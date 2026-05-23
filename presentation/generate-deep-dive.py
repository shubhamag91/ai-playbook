#!/usr/bin/env python3
"""AI Playbook Deep-Dive Presentation Generator — Mathematics, Inference & Systems Engineering."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
import os

# Colors
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF8, 0xF9, 0xFA)  # Light warm gray bg
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
DKGRAY  = RGBColor(0x4A, 0x4A, 0x4A)
MDGRAY  = RGBColor(0x88, 0x88, 0x80)
LTGRAY  = RGBColor(0xE0, 0xDC, 0xD6)
ACCENT  = RGBColor(0x3B, 0x82, 0xF6)  # Blue
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
M = Inches(1.0)  # margin

TOTAL_SLIDES = 14

# ─── Helpers ───────────────────────────────────────────────
def new_slide(bg_color=BG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg_color
    return s

def header_bar(slide, text, color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = M
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return bar

def slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{num} / {TOTAL_SLIDES}"; p.font.size = Pt(9); p.font.color.rgb = MDGRAY; p.alignment = PP_ALIGN.RIGHT

def notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_textbox(slide, left, top, width, height, text, font_size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=14, color=BLACK, spacing=Pt(12), bullet_char="\u25CF"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size); p.font.color.rgb = color
        p.space_after = spacing
    return txBox

def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.45), font_size=11):
    n_rows = len(rows) + 1; n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows).table
    col_w = (int(width) - 300000) / n_cols
    for i in range(n_cols): tbl.columns[i].width = int(col_w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs: p.font.size = Pt(font_size-1); p.font.color.rgb = WHITE; p.font.bold = True
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(val)
            for p in c.text_frame.paragraphs: p.font.size = Pt(font_size); p.font.color.rgb = BLACK
            if i % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = LTGRAY
    return tbl

def make_shape(slide, shape_type, left, top, width, height, fill_color=ACCENT, text="", font_size=10, font_color=WHITE):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    if text:
        tf = shp.text_frame; tf.word_wrap = True; tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(font_size); p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
    return shp

def make_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color; connector.line.width = Pt(2)
    return connector

def metric_card(slide, left, top, width, height, number, label, color=ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LTGRAY; card.line.width = Pt(1)
    tf = card.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.text = str(number); p1.font.size = Pt(24); p1.font.color.rgb = color
    p1.font.bold = True; p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(10); p2.font.color.rgb = MDGRAY
    p2.alignment = PP_ALIGN.CENTER
    return card

def divider_slide(section_num, title, subtitle=""):
    s = new_slide(ACCENT)
    add_textbox(s, M, Inches(2.0), Inches(11.3), Inches(1.5), f"Part {section_num}", 18, WHITE, False, PP_ALIGN.LEFT)
    add_textbox(s, M, Inches(3.0), Inches(11.3), Inches(1.5), title, 36, WHITE, True, PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(s, M, Inches(4.2), Inches(11.3), Inches(1.0), subtitle, 16, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.LEFT)
    return s

# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(1.8), Inches(11.3), Inches(1.5), "AI Research & Deep-Dives: Math & Systems", 40, ACCENT, True, PP_ALIGN.LEFT)
add_textbox(s, M, Inches(3.2), Inches(11.3), Inches(1.0), "Transformer Internals, Reasoning Paradigms, & Inference Optimization", 22, BLACK, False, PP_ALIGN.LEFT)
add_textbox(s, M, Inches(4.2), Inches(11.3), Inches(0.6), "Level 3: Deep-Dive Presentation  \u2022  AI Playbook", 13, MDGRAY, False, PP_ALIGN.LEFT)
make_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(3.1), Inches(1.8), Pt(4), ACCENT)
slide_number(s, 1)
notes(s, "Welcome to the Level 3 Deep-Dive. Today we will bypass conceptual summaries to explore the deep mathematics, GPU kernel optimizations, active scaling bottlenecks, reasoning path search algorithms, and evaluations.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — PART 1 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(1, "Transformer Internals & Scaling Math", "Attention mathematics, GPU kernel fusion, MoE gating, and cache footprints")
slide_number(prs.slides[-1], 2)
notes(prs.slides[-1], "We start with the structural mathematics of modern autoregressive models, exploring why scaled attention exists, how FlashAttention operates on GPUs, and mapping KV cache size limits.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE ATTENTION MECHANISM
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "The Mathematical Core: Scaled Dot-Product Attention")

# Equation Card
eq_card = make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(1.0), Inches(5.3), Inches(2.2), RGBColor(0x1A, 0x2D, 0x4D))
add_textbox(s, Inches(1.2), Inches(1.5), Inches(4.9), Inches(1.2),
    "Attention(Q, K, V) = softmax( QK^T / \u221Ad_k ) V",
    22, WHITE, True, PP_ALIGN.CENTER)

# Explanations
add_textbox(s, M, Inches(3.4), Inches(5.3), Inches(3.0),
    "Where variables denote:\n"
    "\u2022 Q (Query): The token seeking context.\n"
    "\u2022 K (Key): The tokens offering context.\n"
    "\u2022 V (Value): The actual information payload.\n"
    "\u2022 d_k: The dimension of query and key vectors.",
    12.5, BLACK)

# Math explanation card
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.0), Inches(5.5), Inches(4.6), WHITE)
s.shapes[-1].line.color.rgb = ACCENT; s.shapes[-1].line.width = Pt(1.5)

add_textbox(s, Inches(7.0), Inches(1.15), Inches(5.1), Inches(0.4), "Why the scaling factor \u221Ad_k matters:", 15, ACCENT, True)
add_textbox(s, Inches(7.0), Inches(1.6), Inches(5.1), Inches(3.8),
    "As dimension d_k grows large, the dot products grow large in magnitude.\n\n"
    "This pushes the softmax function into regions with extremely small gradients (the vanishing gradient problem).\n\n"
    "Dividing by \u221Ad_k stabilizes the variance of the dot products to 1, ensuring softmax outputs do not saturate at 0 or 1.\n\n"
    "Softmax Saturation = Softmax derivative drops to 0, which halts backpropagation training immediately.",
    12, BLACK)

# callout
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(0.8), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.95), Inches(10.7), Inches(0.5),
    "Softmax saturation makes gradient updates impossible. The scaling factor is the key to training large hidden states.", 11, ACCENT, True)

slide_number(s, 3)
notes(s, "Dot product attention forms the core. Explain that without dividing by the square root of d_k, training deep architectures is physically impossible due to softmax gradient death.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — HARDWARE-LEVEL ENHANCEMENTS
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Hardware Optimizations: RoPE & FlashAttention")

# Column 1: RoPE
col_w = Inches(5.3); col_x = [M, Inches(7.0)]
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[0], Inches(1.0), col_w, Inches(4.6), WHITE)
s.shapes[-1].line.color.rgb = ACCENT; s.shapes[-1].line.width = Pt(1.5)

add_textbox(s, col_x[0] + Inches(0.2), Inches(1.15), col_w - Inches(0.4), Inches(0.4), "Rotary Position Embeddings (RoPE)", 16, ACCENT, True)
add_textbox(s, col_x[0] + Inches(0.2), Inches(1.65), col_w - Inches(0.4), Inches(3.8),
    "Instead of adding absolute position markers to word embeddings, RoPE applies a 2D rotation matrix directly to Q and K projections in complex vector space.\n\n"
    "\u2022 Preserves relative distances: The decay of rotation mathematically aligns with natural token distance.\n"
    "\u2022 Dynamic Extrapolation: Allows models to scale from 8K training context windows to 1M inference context windows at runtime without retraining.",
    11.5, BLACK)

# Column 2: FlashAttention
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[1], Inches(1.0), col_w, Inches(4.6), WHITE)
s.shapes[-1].line.color.rgb = GREEN; s.shapes[-1].line.width = Pt(1.5)

add_textbox(s, col_x[1] + Inches(0.2), Inches(1.15), col_w - Inches(0.4), Inches(0.4), "FlashAttention (Kernel Fusion)", 16, GREEN, True)
add_textbox(s, col_x[1] + Inches(0.2), Inches(1.65), col_w - Inches(0.4), Inches(3.8),
    "Standard attention is memory-bound (reading/writing massive intermediate matrices from slow High Bandwidth Memory [HBM] to fast GPU SRAM).\n\n"
    "\u2022 FlashAttention fuses the attention kernel to operate entirely in SRAM, bypassing expensive DRAM round-trips.\n"
    "\u2022 Computes softmax increments progressively online rather than storing the full N x N matrix.\n"
    "\u2022 Yields a 2x-4x speedup in wall-clock execution time.",
    11.5, BLACK)

# Takeaway note
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(0.8), RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(s, Inches(1.3), Inches(5.95), Inches(10.7), Inches(0.55),
    "Key: RoPE enables flexible context scaling. FlashAttention bypasses memory transfer bottlenecks on the GPU.", 11, GREEN, True)

slide_number(s, 4)
notes(s, "Detail the GPU constraint. LLMs are memory-bandwidth bound rather than compute-bound. FlashAttention is a software engineering trick that speeds up training and inference purely by optimizing GPU memory hierarchy usage.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — MoE ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Scaling Sparsity: Mixture of Experts (MoE)")

add_textbox(s, M, Inches(1.0), Inches(11.3), Inches(1.2),
    "Mixture-of-Experts (MoE) scales model parameters while maintaining a flat compute footprint. Instead of routing a token through every layer's neural connections, a router directs it to only a small subset of experts.",
    12.5, BLACK)

# MoE Flow Diagram
y = Inches(2.2); w_box = Inches(2.1); h_box = Inches(1.4); gap = Inches(0.6)
experts = [
    ("Router Network\n\nAnalyzes token & assigns weights", ACCENT),
    ("Expert 1 (Math)\n\nProcesses reasoning tokens", DKGRAY),
    ("Expert 2 (Code)\n\nProcesses syntax tokens", DKGRAY),
    ("Expert 3 (Writing)\n\nProcesses prose tokens", DKGRAY),
    ("Fused Output\n\nCombines top 2 expert outputs", GREEN)
]

for i, (label, color) in enumerate(experts):
    # Position
    if i == 0:
        x = M; y_pos = y + Inches(0.7)
    elif i == 4:
        x = M + Inches(9.2); y_pos = y + Inches(0.7)
    else:
        x = M + Inches(4.5); y_pos = y + (i-1) * Inches(1.0)
        
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, w_box, h_box if i in [0,4] else Inches(0.85), color, label, 10.5, WHITE)
    
    # Arrows
    if i == 0:
        for j in range(3):
            make_arrow(s, x + w_box, y_pos + h_box/2, M + Inches(4.5), y + j * Inches(1.0) + Inches(0.42), ACCENT)
    elif i in [1, 2, 3]:
        make_arrow(s, x + w_box, y_pos + Inches(0.42), M + Inches(9.2), y + Inches(0.7) + h_box/2, ACCENT)

# Formula breakdown
add_textbox(s, M, Inches(5.3), Inches(11.3), Inches(1.4),
    "Mathematical routing breakdown:\n"
    "\u2022 Total Parameters vs. Active Parameters: Llama 4 MoE might host 405B parameters, but only routes to 40B active parameters per token. Gating router applies G(x) softmax across experts to select Top-K experts dynamically.",
    12, BLACK)

slide_number(s, 5)
notes(s, "Explain MoE routing. The gating network decides which expert network is best suited for each token. This allows models to act highly knowledgeable on specialized topics without inflating active inference speed.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — KV CACHE SCALING MATH
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Active Memory Bottleneck: KV Cache Scaling Math")

# Formula Box
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(1.0), Inches(11.3), Inches(1.1), RGBColor(0x1A, 0x2D, 0x4D))
add_textbox(s, Inches(1.2), Inches(1.15), Inches(10.9), Inches(0.85),
    "KV Cache Footprint (Bytes) = 2 \u00D7 2 \u00D7 n_layers \u00D7 n_heads \u00D7 d_head \u00D7 s \u00D7 b",
    20, WHITE, True, PP_ALIGN.CENTER)

# Details
add_textbox(s, M, Inches(2.3), Inches(5.3), Inches(3.2),
    "Where equations variables denote:\n"
    "\u2022 n_layers: Number of transformer decoder layers.\n"
    "\u2022 n_heads: Number of attention key-value heads.\n"
    "\u2022 d_head: Dimension of each head vector.\n"
    "\u2022 s: Input sequence length (context depth).\n"
    "\u2022 b: Inference batch size (parallel users).\n\n"
    "Note: First factor of 2 accounts for Keys & Values; second factor accounts for FP16/BF16 data formats (2 bytes).",
    12, BLACK)

# Example Table
headers = ["Model", "Layers", "Heads (KV)", "Head Dim", "Sequence", "Batch", "KV Cache Size"]
rows = [
    ["Llama 3 8B", "32", "8", "128", "8,192", "1", "0.50 GB"],
    ["Llama 3 8B", "32", "8", "128", "8,192", "32", "16.00 GB"],
    ["Llama 3 70B", "80", "8", "128", "8,192", "1", "1.31 GB"],
    ["Llama 3 70B", "80", "8", "128", "128,000", "1", "20.48 GB"]
]
add_table(s, headers, rows, Inches(6.6), Inches(2.3), Inches(5.7), Inches(0.48), 10)

# Bottom card
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.7), Inches(11.3), Inches(0.9), RGBColor(0xFE, 0xF2, 0xF2))
add_textbox(s, Inches(1.3), Inches(5.85), Inches(10.7), Inches(0.6),
    "Strategic Bottleneck: KV Cache grows linearly with sequence length and batch size. Multi-Query Attention (MQA) and Grouped-Query Attention (GQA) reduce KV heads to mitigate VRAM overflow.",
    11, RED, True)

slide_number(s, 6)
notes(s, "Explain the linear scaling problem. Because we must store keys and values in GPU RAM to avoid recalculation, longer prompts eat up massive amounts of memory. Explain GQA as the primary architecture solution.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — PART 2 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(2, "Reasoning Models & Test-Time Compute", "Algorithmic path searches, GRPO reinforcement learning, and outcome vs. process rewards")
slide_number(prs.slides[-1], 7)
notes(prs.slides[-1], "In Part 2, we explore the paradigm shift to reasoning systems like o1/o3 and DeepSeek-R1, and inspect the algorithms that govern thinking paths.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — REASONING PARADIGMS
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "The Paradigm Shift: Pretraining vs. Test-Time Compute")

# Comparison layout
col_w = Inches(5.3); col_x = [M, Inches(7.0)]

make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[0], Inches(1.1), col_w, Inches(4.5), WHITE)
s.shapes[-1].line.color.rgb = ACCENT; s.shapes[-1].line.width = Pt(1.5)

add_textbox(s, col_x[0] + Inches(0.2), Inches(1.25), col_w - Inches(0.4), Inches(0.4), "Standard Models (Pretraining Compute)", 15, ACCENT, True)
add_textbox(s, col_x[0] + Inches(0.2), Inches(1.75), col_w - Inches(0.4), Inches(3.6),
    "In standard models, intelligence is proportional to parameter size and training tokens.\n\n"
    "\u2022 Flattens at inference: Model spits out the next token using a static, fixed forward-pass compute budget.\n"
    "\u2022 Cannot allocate extra compute: An easy 'Yes/No' takes the exact same forward-pass parameters as a complex math proof.\n"
    "\u2022 High training costs, fixed query latency.",
    11.5, BLACK)

make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[1], Inches(1.1), col_w, Inches(4.5), WHITE)
s.shapes[-1].line.color.rgb = GREEN; s.shapes[-1].line.width = Pt(1.5)

add_textbox(s, col_x[1] + Inches(0.2), Inches(1.25), col_w - Inches(0.4), Inches(0.4), "Reasoning Models (Test-Time Compute)", 15, GREEN, True)
add_textbox(s, col_x[1] + Inches(0.2), Inches(1.75), col_w - Inches(0.4), Inches(3.6),
    "Reasoning systems scale compute dynamically during generation using search algorithms.\n\n"
    "\u2022 Algorithmic paths: Evaluates multiple possible thoughts, backtracks when finding logical dead-ends, and self-corrects.\n"
    "\u2022 Thinking tokens: Allocates hundreds of tokens internally to 'think' before producing a final output.\n"
    "\u2022 Scales compute at inference: Harder questions generate longer execution paths to yield superior outcomes.",
    11.5, BLACK)

# Bottom note
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(0.8), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.95), Inches(10.7), Inches(0.55),
    "Key: Intelligence is no longer frozen in SFT weights. We can scale performance purely by allowing longer inference thinking paths.", 11, ACCENT, True)

slide_number(s, 8)
notes(s, "Explain that reasoning models decouple model size from capability. An 8B parameter model trained with test-time search can beat a massive 70B static model on complex math and coding tasks.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — SEARCH & REWARD ALGORITHMS
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Search & Reward: CoT, ToT, and Process Reward Models")

# Three columns
col_w = Inches(3.5); col_x = [M, Inches(4.9), Inches(8.8)]

# CoT
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[0], Inches(1.1), col_w, Inches(4.4), WHITE)
s.shapes[-1].line.color.rgb = ACCENT; s.shapes[-1].line.width = Pt(1.5)
add_textbox(s, col_x[0] + Inches(0.15), Inches(1.25), col_w - Inches(0.3), Inches(0.45), "Chain-of-Thought (CoT)", 14, ACCENT, True)
add_textbox(s, col_x[0] + Inches(0.15), Inches(1.8), col_w - Inches(0.3), Inches(3.6),
    "Forces a linear reasoning path.\n\n"
    "By outputting step-by-step logic ('Let's think step by step') before the final answer, standard models solve complex math.\n\n"
    "Math grounding: CoT increases target prediction accuracy by allowing intermediate computations to sit in the context window history.",
    11, BLACK)

# ToT
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[1], Inches(1.1), col_w, Inches(4.4), WHITE)
s.shapes[-1].line.color.rgb = GREEN; s.shapes[-1].line.width = Pt(1.5)
add_textbox(s, col_x[1] + Inches(0.15), Inches(1.25), col_w - Inches(0.3), Inches(0.45), "Tree-of-Thoughts (ToT)", 14, GREEN, True)
add_textbox(s, col_x[1] + Inches(0.15), Inches(1.8), col_w - Inches(0.3), Inches(3.6),
    "Generalizes CoT into tree branches.\n\n"
    "Evaluates multiple parallel reasoning paths.\n\n"
    "Uses search algorithms like Beam Search or Depth-First Search (DFS) to explore logical branches and prune paths that score poorly, avoiding cognitive dead-ends.",
    11, BLACK)

# PRMs
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[2], Inches(1.1), col_w, Inches(4.4), WHITE)
s.shapes[-1].line.color.rgb = AMBER; s.shapes[-1].line.width = Pt(1.5)
add_textbox(s, col_x[2] + Inches(0.15), Inches(1.25), col_w - Inches(0.3), Inches(0.45), "Process Reward Models (PRM)", 14, AMBER, True)
add_textbox(s, col_x[2] + Inches(0.15), Inches(1.8), col_w - Inches(0.3), Inches(3.6),
    "Traditional Outcome Models (ORMs) only grade the final answer.\n\n"
    "PRMs grade every single step in the reasoning chain.\n\n"
    "Critical benefit: Provides dense feedback to guide searches, detects logic bugs early, and mitigates 'accidental' correct answers derived from incorrect calculations.",
    11, BLACK)

slide_number(s, 9)
notes(s, "Differentiate these methods. CoT is linear. ToT is branching with search. PRM is the judge grading each logic step. DeepSeek-R1 combines GRPO reinforcement learning with PRM-like metrics to align reasoning outputs.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — PART 3 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(3, "Advanced Inference & Evaluation", "Speculative decoding, quantization limits, RAGAS evaluations, and system designs")
slide_number(prs.slides[-1], 10)
notes(prs.slides[-1], "In the final section, we look at deployment realities: speculative decoding, precision quantization tradeoffs, the mathematics of RAGAS, and custom routing pipelines.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — INFERENCE OPTIMIZATIONS
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Advanced Optimizations: Speculative Decoding & Quantization")

col_w = Inches(5.3); col_x = [M, Inches(7.0)]

# Speculative Decoding
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[0], Inches(1.1), col_w, Inches(4.5), WHITE)
s.shapes[-1].line.color.rgb = ACCENT; s.shapes[-1].line.width = Pt(1.5)

add_textbox(s, col_x[0] + Inches(0.2), Inches(1.25), col_w - Inches(0.4), Inches(0.4), "Speculative Decoding", 15, ACCENT, True)
add_textbox(s, col_x[0] + Inches(0.2), Inches(1.75), col_w - Inches(0.4), Inches(3.6),
    "Generates tokens using a lightweight 'draft' model (e.g., 8B) very quickly, then verifies them in parallel using a large 'target' model (e.g., 70B) in a single GPU pass.\n\n"
    "\u2022 Parallel validation: Because validating tokens is compute-bound, the large model validates 5-10 tokens simultaneously.\n"
    "\u2022 Yields a 2x-3x speedup with mathematically identical quality outputs.",
    11.5, BLACK)

# Quantization
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[1], Inches(1.1), col_w, Inches(4.5), WHITE)
s.shapes[-1].line.color.rgb = GREEN; s.shapes[-1].line.width = Pt(1.5)

add_textbox(s, col_x[1] + Inches(0.2), Inches(1.25), col_w - Inches(0.4), Inches(0.4), "Model Quantization (FP8, INT8, FP4)", 15, GREEN, True)
add_textbox(s, col_x[1] + Inches(0.2), Inches(1.75), col_w - Inches(0.4), Inches(3.6),
    "Compresses model parameters by converting floating-point weights (FP32/FP16) into lower-precision formats like FP8 or INT8.\n\n"
    "\u2022 VRAM Reduction: FP16 to FP8 cuts memory footprint in half. A 70B model requires ~140GB in FP16 but fits in ~70GB in FP8.\n"
    "\u2022 Minimal Quality Loss: FP8 and INT8 conversions maintain near-baseline accuracy on benchmark evaluations.",
    11.5, BLACK)

slide_number(s, 11)
notes(s, "Speculative decoding and quantization are the twin workhorses of LLMOps. Explain how they make massive models accessible on standard GPU servers.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — RAGAS EVALUATION FRAMEWORK
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "RAG Evaluation: The RAGAS Mathematical Framework")

# The four metrics
metrics = [
    ("Faithfulness", "Is the answer grounded in the retrieved context?",
     "Faithfulness = Factual Sentences in Answer / Total Sentences in Answer", ACCENT),
     
    ("Answer Relevance", "Does the answer directly address the user's prompt?",
     "Answer Relevance = Cosine Similarity of Generated Questions to Query", GREEN),
     
    ("Context Recall", "Did the system retrieve all the necessary facts?",
     "Context Recall = Sentences in Context Ground Truth / Total Ground Truth", AMBER),
     
    ("Context Precision", "Is the retrieved context clean and noise-free?",
     "Context Precision = Precision@K over retrieved document ranks", RGBColor(0x1A, 0x2D, 0x4D))
]

for i, (name, desc, formula, color) in enumerate(metrics):
    col = i % 2; row = i // 2
    x = M + col * Inches(5.8); y = Inches(1.1) + row * Inches(2.2)
    
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(1.9), WHITE)
    s.shapes[-1].line.color.rgb = color; s.shapes[-1].line.width = Pt(1.5)
    
    add_textbox(s, x + Inches(0.2), y + Inches(0.12), Inches(5.1), Inches(0.3), name, 13.5, color, True)
    add_textbox(s, x + Inches(0.2), y + Inches(0.45), Inches(5.1), Inches(0.4), desc, 10, MDGRAY)
    
    # Formula box inside card
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.95), Inches(5.1), Inches(0.75), RGBColor(0xFA, 0xFA, 0xF9))
    s.shapes[-1].line.color.rgb = LTGRAY; s.shapes[-1].line.width = Pt(1)
    
    add_textbox(s, x + Inches(0.3), y + Inches(1.1), Inches(4.9), Inches(0.55), formula, 9.5, DKGRAY, True)

slide_number(s, 12)
notes(s, "Explain the RAGAS framework. Standard software testing fails with LLMs due to non-deterministic outputs. RAGAS uses LLMs as judges to score faithfulness, relevance, recall, and precision mathematically.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — ADVANCED SYSTEM DESIGN PATTERNS
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "System Design Patterns: Cascade Difficulty Routing")

add_textbox(s, M, Inches(1.0), Inches(11.3), Inches(1.2),
    "Instead of routing every request to a massive, expensive model, Cascade Routing employs an intelligent classifier to triage tasks by complexity. Simple requests are resolved by cheap local models, while hard logic is escalated.",
    12.5, BLACK)

# Routing Diagram
y = Inches(2.2)
# Query
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y + Inches(1.0), Inches(2.2), Inches(0.85), RGBColor(0x1A, 0x2D, 0x4D), "User Prompt Input", 11, WHITE)
make_arrow(s, Inches(3.2), y + Inches(1.42), Inches(4.5), y + Inches(1.42), ACCENT)

# Classifier
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), y + Inches(0.65), Inches(2.3), Inches(1.5), ACCENT, "Intelligent Classifier\n\n(Classifies intent &\ncomplexity score)", 11, WHITE)

# Escalation paths
make_arrow(s, Inches(6.8), y + Inches(1.0), Inches(8.5), y + Inches(0.5), GREEN)
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), y, Inches(3.5), Inches(0.85), GREEN, "Easy / Fast Path\n\nRoute to Llama 8B / DeepSeek\n(Latency: ~100ms  |  Cost: $0)", 9.5, WHITE)

make_arrow(s, Inches(6.8), y + Inches(1.8), Inches(8.5), y + Inches(2.3), RED)
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), y + Inches(1.9), Inches(3.5), Inches(0.85), RED, "Hard / Reasoning Path\n\nRoute to o3 / DeepSeek-R1\n(Latency: 5-15s  |  Cost: $15/1M)", 9.5, WHITE)

# Details
add_textbox(s, M, Inches(5.1), Inches(11.3), Inches(1.3),
    "Design Advantages:\n"
    "\u2022 Yields up to 90% cost savings by preventing wasteful allocations on simple greeting or search queries.\n"
    "\u2022 Maximizes system availability and keeps average user latency to a minimum.",
    11.5, BLACK)

slide_number(s, 13)
notes(s, "Explain Cascade Routing. It is a critical design pattern in LLMOps. By categorizing prompts early, we achieve the latency and cost footprint of a small model with the capabilities of a frontier system.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(1.8), Inches(11.3), Inches(1.0), "Deep-Dive Q&A", 52, ACCENT, True, PP_ALIGN.CENTER)
add_textbox(s, M, Inches(3.0), Inches(11.3), Inches(0.8), "Ask me anything about scaled attention, KV caching constraints, thinking paths, or speculative optimizations", 18, BLACK, False, PP_ALIGN.CENTER)
add_textbox(s, M, Inches(4.5), Inches(11.3), Inches(0.6), "Visit the Playbook: ai-playbook.pages.dev  \u2022  FlashAttention  \u2022  vLLM", 14, MDGRAY, False, PP_ALIGN.CENTER)
slide_number(s, 14)
notes(s, "Open floor for research and deep-dive technical questions. Address topics such as flash-attention hardware constraints, RoPE sequence extrapolation limits, or reinforcement learning GRPO implementations.")

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ai-playbook-deep-dive.pptx")
prs.save(output_path)
print(f"\u2713 Saved {output_path}")
