#!/usr/bin/env python3
"""AI Playbook presentation v2 — Corporate light theme, charts, diagrams."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import os

# Colors
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF8, 0xF9, 0xFA)  # light gray bg
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
DKGRAY  = RGBColor(0x4A, 0x4A, 0x4A)
MDGRAY  = RGBColor(0x88, 0x88, 0x80)
LTGRAY  = RGBColor(0xE0, 0xDC, 0xD6)
ACCENT  = RGBColor(0x3B, 0x82, 0xF6)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
M = Inches(1.0)  # margin

# ─── Helpers ───────────────────────────────────────────────
def new_slide(bg_color=BG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg_color
    return s

def header_bar(slide, text, color=ACCENT):
    """Full-width accent bar at top with section text."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = M
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return bar

def slide_number(slide, num, total=23):
    txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = str(num); p.font.size = Pt(9); p.font.color.rgb = MDGRAY; p.alignment = PP_ALIGN.RIGHT

def notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_textbox(slide, left, top, width, height, text, font_size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=14, color=BLACK, spacing=Pt(12), bullet_char="\u25CF"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size); p.font.color.rgb = color
        p.space_after = spacing
    return txBox

def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.45), font_size=11):
    n_rows = len(rows) + 1; n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows).table
    col_w = (int(width) - 300000) / n_cols
    for i in range(n_cols): tbl.columns[i].width = int(col_w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs: p.font.size = Pt(font_size-1); p.font.color.rgb = WHITE; p.font.bold = True
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(val)
            for p in c.text_frame.paragraphs: p.font.size = Pt(font_size); p.font.color.rgb = BLACK
            if i % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = LTGRAY
    return tbl

def color_cell(table, row, col, bg):
    c = table.cell(row, col)
    c.fill.solid(); c.fill.fore_color.rgb = bg

def make_shape(slide, shape_type, left, top, width, height, fill_color=ACCENT, text="", font_size=10, font_color=WHITE):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    if text:
        tf = shp.text_frame; tf.word_wrap = True; tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(font_size); p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
    return shp

def make_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color; connector.line.width = Pt(2)
    return connector

def metric_card(slide, left, top, width, height, number, label, color=ACCENT):
    """A highlighted stat card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LTGRAY; card.line.width = Pt(1)
    tf = card.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.text = str(number); p1.font.size = Pt(24); p1.font.color.rgb = color
    p1.font.bold = True; p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(10); p2.font.color.rgb = MDGRAY
    p2.alignment = PP_ALIGN.CENTER
    return card

def divider_slide(section_num, title, subtitle=""):
    s = new_slide(ACCENT)
    add_textbox(s, M, Inches(2.0), Inches(11.3), Inches(1.5), f"Part {section_num}", 18, WHITE, False, PP_ALIGN.LEFT)
    add_textbox(s, M, Inches(3.0), Inches(11.3), Inches(1.5), title, 36, WHITE, True, PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(s, M, Inches(4.2), Inches(11.3), Inches(1.0), subtitle, 16, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.LEFT)
    return s


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(1.5), Inches(11.3), Inches(1.5), "AI & LLMs", 48, ACCENT, True, PP_ALIGN.LEFT)
add_textbox(s, M, Inches(3.0), Inches(11.3), Inches(1.0), "A Practical Introduction for Decision Makers", 24, BLACK, False, PP_ALIGN.LEFT)
add_textbox(s, M, Inches(4.0), Inches(11.3), Inches(0.6), "From the AI Playbook — May 2026", 14, MDGRAY, False, PP_ALIGN.LEFT)
# Blue decorative line
make_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(3.0), Inches(1.5), Pt(4), ACCENT)
slide_number(s, 1)
notes(s, "Welcome. Today's goal: a practical understanding of AI/LLMs — what they are, how they work, which tools matter, and how to get started. No prior AI knowledge needed.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — PART 1 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(1, "The AI Landscape", "What AI is, where we are, and what's real vs hype")
slide_number(prs.slides[-1], 2)
notes(prs.slides[-1], "We'll start with the big picture — what LLMs are, the current landscape, and separate reality from hype.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT ARE LLMs? (+ TIMELINE)
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "What Are LLMs?")
# Left: definition
add_textbox(s, M, Inches(1.0), Inches(5.5), Inches(2.5),
    'A Large Language Model is software that predicts the next word in a sentence.\n\n'
    'Think of autocomplete on your phone — but trained on trillions of words from the internet, books, and code.\n\n'
    'Not magic. Pattern matching at massive scale.\n\n'
    '"Large" refers to three things:\n'
    '• Training data: trillions of words\n'
    '• Parameters: billions of prediction "dials"\n'
    '• Context: up to 1 million tokens of memory',
    14, BLACK)
# Right: timeline
add_textbox(s, Inches(7.2), Inches(1.0), Inches(5.0), Inches(0.5), "Evolution (2017\u20132026)", 16, ACCENT, True)
timeline_events = [
    (Inches(7.2), Inches(1.8), "2017", "Transformers paper"),
    (Inches(7.2), Inches(2.4), "2020", "GPT-3 (175B)"),
    (Inches(7.2), Inches(3.0), "2022", "ChatGPT launch"),
    (Inches(7.2), Inches(3.6), "2023", "GPT-4, Claude 2"),
    (Inches(7.2), Inches(4.2), "2024", "Claude 3.5, GPT-4o"),
    (Inches(7.2), Inches(4.8), "2025", "o1, o3 reasoning"),
    (Inches(7.2), Inches(5.4), "2026", "Claude 4.7, GPT-5.5"),
]
# Vertical line
make_shape(s, MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(1.8), Pt(2), Inches(4.0), ACCENT)
for x, y, yr, txt in timeline_events:
    make_shape(s, MSO_SHAPE.OVAL, Inches(7.92), y, Inches(0.15), Inches(0.15), ACCENT)
    add_textbox(s, Inches(8.3), y - Pt(2), Inches(0.7), Inches(0.3), yr, 11, ACCENT, True)
    add_textbox(s, Inches(9.0), y - Pt(2), Inches(3.0), Inches(0.3), txt, 10, BLACK)
# Callout at bottom
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(0.9), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.95), Inches(10.7), Inches(0.6),
    '\u25B6 Key: Prices dropping 50%+ per year. What cost $1 in 2023 costs ~$0.10 today. Frontier models are becoming free.',
    12, ACCENT, True)
slide_number(s, 3)
notes(s, "LLM = software that predicts text. The word 'large' = lots of data, lots of parameters, lots of context. Timeline shows astonishing pace: 2017 transformer paper to fully capable reasoning models in just 9 years.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — 2026 LANDSCAPE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "The 2026 AI Landscape")
# Metric cards
metric_card(s, M, Inches(1.0), Inches(2.5), Inches(1.0), "50%+", "Price drop per year", GREEN)
metric_card(s, Inches(3.8), Inches(1.0), Inches(2.5), Inches(1.0), "1M", "Max context tokens", ACCENT)
metric_card(s, Inches(6.6), Inches(1.0), Inches(2.5), Inches(1.0), "87%", "Devs using AI tools", AMBER)
metric_card(s, Inches(9.4), Inches(1.0), Inches(2.8), Inches(1.0), "$0.14/1M", "Cheapest frontier model", GREEN)
# Key players table
add_table(s, ["Company", "Model", "Context", "Best Strength", "Input Price/1M"],
    [["Anthropic", "Claude Opus 4.7", "400K", "Best reasoning & writing", "$15"],
     ["OpenAI", "GPT-5.5", "128K", "All-purpose, fastest", "$2"],
     ["Google", "Gemini 3.1 Pro", "1M (largest)", "Research, long docs", "$2"],
     ["DeepSeek", "DeepSeek V4", "256K", "10-50x cheaper, open", "$0.14"],
     ["Meta", "Llama 4 Scout", "128K", "Open-weight, free", "Free (MIT)"],
     ["xAI", "Grok 3 Pro", "256K", "Real-time X/Twitter data", "API pricing"]],
    M, Inches(2.3), Inches(11.3), Inches(0.47))
# Trend callout
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.3), Inches(11.3), Inches(1.4), RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(s, Inches(1.3), Inches(5.4), Inches(10.7), Inches(1.2),
    'Key Trends\n'
    '• Open models (Llama, DeepSeek, Qwen) matching closed models on benchmarks\n'
    '• Context windows growing from 8K (2022) → 400K (2025) → 1M tokens (2026)\n'
    '• Agentic AI moving from demos to production — tools that plan, code, and execute',
    12, GREEN, False)
slide_number(s, 4)
notes(s, "The landscape at May 2026. Four stats capture the pace of change. DeepSeek at $0.14/1M tokens is 100x cheaper than Claude Opus at $15/1M. Open-weight means free to self-host. This changes the economics of building with AI.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODEL COMPARISON (CAPABILITY HEATMAP)
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Model Capability Comparison")
add_table(s, ["Model", "Writing", "Coding", "Reasoning", "Context", "Cost/Value"],
    [["Claude Opus 4.7",  "\u2605\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605",  "\u2605\u2605\u2605\u2605\u2605", "400K", "\u2605\u2605"],
     ["GPT-5.5",         "\u2605\u2605\u2605\u2605",   "\u2605\u2605\u2605\u2605",  "\u2605\u2605\u2605\u2605",   "128K", "\u2605\u2605\u2605\u2605"],
     ["Gemini 3.1 Pro",  "\u2605\u2605\u2605\u2605",   "\u2605\u2605\u2605",     "\u2605\u2605\u2605\u2605",   "1M",   "\u2605\u2605\u2605\u2605"],
     ["DeepSeek V4",     "\u2605\u2605\u2605",        "\u2605\u2605\u2605\u2605",  "\u2605\u2605\u2605\u2605\u2605", "256K", "\u2605\u2605\u2605\u2605\u2605"],
     ["Llama 4 Scout",   "\u2605\u2605\u2605",        "\u2605\u2605\u2605\u2605",  "\u2605\u2605\u2605",       "128K", "\u2605\u2605\u2605\u2605\u2605"],
     ["Grok 3 Pro",      "\u2605\u2605\u2605\u2605",   "\u2605\u2605\u2605\u2605",  "\u2605\u2605\u2605\u2605",   "256K", "\u2605\u2605\u2605"]],
    M, Inches(1.0), Inches(11.3), Inches(0.65))
# Key
add_textbox(s, M, Inches(5.2), Inches(11.3), Inches(0.5),
    '\u2605 = Star rating: 5 = world-class, 1 = limited. Based on published benchmarks and public evaluations.',
    10, MDGRAY)
# Bottom callout
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(1.0), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.9), Inches(10.7), Inches(0.8),
    'Key Insight: There is no single "best" model. Claude excels at reasoning. GPT at speed. Gemini at context. DeepSeek at cost. Llama at control. Choose based on your specific use case.',
    12, ACCENT)
slide_number(s, 5)
notes(s, "Star ratings show relative capability. DeepSeek matches or exceeds Claude on reasoning at 100x lower cost. Llama is free but 2 stars on reasoning. The gap between open and closed is closing fast.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — REALITY VS HYPE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Reality Check: What AI Can & Cannot Do")

# Left: Myths
add_textbox(s, M, Inches(0.9), Inches(5.5), Inches(0.4), "Common Myths", 16, RED, True)
myths = [
    '"LLMs understand language" — They pattern-match. Understanding is our projection.',
    '"LLMs are general intelligences" — They\'re narrow: good at text, weak at uncertainty.',
    '"ChatGPT is always right" — No. They hallucinate. Verify critical facts.',
    '"More parameters = smarter" — Architecture and data quality matter more.',
    '"LLMs will replace humans" — Humans + LLMs outperform either alone.',
]
add_bullets(s, M, Inches(1.5), Inches(5.5), Inches(4.5), myths, 12, BLACK, Pt(14), "\u2715")

# Right: Reality
add_textbox(s, Inches(7.2), Inches(0.9), Inches(5.5), Inches(0.4), "What They Actually Do", 16, GREEN, True)
reality = [
    'Write & edit text — often better than most people.',
    'Explain complex topics at any level of detail.',
    'Code generation & review — saves 40-60% dev time.',
    'Summarize long documents into actionable insights.',
    'Brainstorm & ideate — tireless creative partner.',
    'Translate between languages and formats (JSON, SQL, etc).',
]
add_bullets(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(4.5), reality, 12, BLACK, Pt(12), "\u2713")

# Bottom: Limitations
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.5), Inches(11.3), Inches(1.2), RGBColor(0xFE, 0xF2, 0xF2))
add_textbox(s, Inches(1.3), Inches(5.6), Inches(10.7), Inches(1.0),
    'Limitations: No internet access (trained not connected) • Training cutoff date matters • Hallucination risk • Not domain experts\n'
    'Solution: RAG (feed it your docs), fact-check critical outputs, use AI as a starting point, not the final answer.',
    11, RED)
slide_number(s, 6)
notes(s, "The capabilities slide is mission-critical for setting expectations. Key message: AI is a powerful tool, not magic. Know what it can and can't do.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — PART 2 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(2, "How LLMs Work", "The mechanics: tokens, training, RAG, and prompting")
slide_number(prs.slides[-1], 7)
notes(prs.slides[-1], "Now we understand what AI can do. Let's understand how it actually works — just enough depth to make informed decisions.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — TRAINING PIPELINE DIAGRAM
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "How Models Learn — The Training Pipeline")

# Build diagram with shapes
y = Inches(1.6); box_w = Inches(2.2); box_h = Inches(1.2); gap = Inches(0.3)
stages = [
    ("Trillions of\nWords\n(internet, books,\ncode)", RGBColor(0x1A, 0x2D, 0x4D)),  # dark blue
    ("Pretraining\n\nLearn language\npatterns & facts", ACCENT),
    ("Fine-tuning\n\nAdapt to specific\ntask or domain", RGBColor(0x60, 0xA5, 0xFA)),
    ("RLHF\n\nHumans rate;\nmodel learns\n\"good vs bad\"", GREEN),
    ("Deployed\nModel\n\nReady for\nAPI / users", RGBColor(0x16, 0x6B, 0x34)),
]
for i, (label, color) in enumerate(stages):
    x = M + i * (box_w + gap + Inches(0.5))
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h, color, label, 11, WHITE)
    if i < len(stages) - 1:
        # Arrow between boxes
        make_arrow(s, x + box_w, y + box_h/2, x + box_w + Inches(0.5), y + box_h/2, ACCENT)

# Description below each stage
y2 = y + box_h + Inches(0.2)
descriptions = [
    ("Data Collection", "Curating training data\nis the hardest part"),
    ("Unsupervised Learning", "Cost: $10M-$100M\nTakes weeks/months"),
    ("Supervised Learning", "Cost: $100-$10K\nTakes hours/days"),
    ("Human Feedback", "Makes models helpful\n& polite — TRAINED, not magic"),
    ("Production", "API or self-hosted\nInference at scale"),
]
for i, (title, desc) in enumerate(descriptions):
    x = M + i * (box_w + gap + Inches(0.5))
    add_textbox(s, x, y2, box_w, Inches(0.65), title, 11, ACCENT, True, PP_ALIGN.CENTER)
    add_textbox(s, x, y2 + Inches(0.25), box_w, Inches(0.5), desc, 9, MDGRAY, False, PP_ALIGN.CENTER)

# Bottom: key takeaway
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(1.0), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.9), Inches(10.7), Inches(0.8),
    'Key: The biggest models cost millions to train. Fine-tuning is much cheaper ($100-$10K). RLHF is why models are polite & helpful — they were TRAINED to be.',
    12, ACCENT, True)
slide_number(s, 8)
notes(s, "Training pipeline: pretraining (expensive, foundation) → fine-tuning (cheap, adapt) → RLHF (makes it polite). Think of pretraining as elementary school, fine-tuning as college major, RLHF as etiquette training.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — TOKENS, PARAMETERS, CONTEXT
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Key Mechanics: Tokens, Parameters & Context")

# Three columns with visuals
col_w = Inches(3.5); col_x = [M, Inches(4.9), Inches(8.8)]

# Column 1: Tokens
add_textbox(s, col_x[0], Inches(1.0), col_w, Inches(0.4), "Tokens", 18, ACCENT, True)
add_textbox(s, col_x[0], Inches(1.5), col_w, Inches(2.5),
    'Models process tokens — not letters, not words.\n\n'
    'A token \u2248 \u00BE of a word.\n'
    '"The quick brown fox" = 4 tokens.\n'
    '~1,000 tokens = ~750 words = ~1.5 pages.\n\n'
    'Pricing is per-token.\n'
    'Input: what the model reads.\n'
    'Output: what the model writes.',
    12, BLACK)
# Token example
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[0], Inches(4.2), col_w, Inches(0.7), LTGRAY)
add_textbox(s, col_x[0], Inches(4.25), col_w, Inches(0.5), "The | quick | brown | fox | jumps", 11, BLACK, False, PP_ALIGN.CENTER)

# Column 2: Parameters
add_textbox(s, col_x[1], Inches(1.0), col_w, Inches(0.4), "Parameters", 18, ACCENT, True)
add_textbox(s, col_x[1], Inches(1.5), col_w, Inches(2.5),
    'Think of parameters as millions of\ntiny "prediction dials."\n\n'
    'GPT-5.5: ~1.76 trillion parameters.\n'
    'Llama 4 405B: 405 billion (open-weight).\n\n'
    'More parameters = more nuance.\n'
    'But architecture & training data\nquality matter equally.',
    12, BLACK)
# Dial visualization
make_shape(s, MSO_SHAPE.OVAL, col_x[1] + Inches(0.5), Inches(3.8), Inches(0.6), Inches(0.6), RGBColor(0x1A, 0x2D, 0x4D))
add_textbox(s, col_x[1] + Inches(1.3), Inches(3.95), col_w - Inches(1.3), Inches(0.4), "\u2248 1.76T dials", 12, ACCENT, True)

# Column 3: Context Windows
add_textbox(s, col_x[2], Inches(1.0), col_w, Inches(0.4), "Context Windows", 18, ACCENT, True)
add_textbox(s, col_x[2], Inches(1.5), col_w, Inches(2.5),
    'Everything the model "remembers"\nin one conversation.\n\n'
    'Gemini: 1M tokens (War & Peace 2x)\n'
    'Claude: 400K (~600 pages)\n'
    'GPT-5.5: 128K (~200 pages)\n\n'
    'Entire conversation sent with\neach message. Long chats cost more.',
    12, BLACK)
# Context bar visualization
add_textbox(s, col_x[2], Inches(3.7), col_w, Inches(0.3), "4K \u2192 128K \u2192 400K \u2192 1M tokens", 9, MDGRAY, False, PP_ALIGN.CENTER)
bar_colors = [MDGRAY, LTGRAY, ACCENT, GREEN]
bar_widths = [Inches(0.4), Inches(0.8), Inches(1.2), Inches(2.0)]
bar_x = col_x[2]
for i in range(4):
    make_shape(s, MSO_SHAPE.RECTANGLE, bar_x, Inches(3.95), bar_widths[i], Inches(0.2), bar_colors[i])
    bar_x += bar_widths[i]
add_textbox(s, col_x[2] + bar_widths[2], Inches(4.2), Inches(2.0), Inches(0.3), "1M \u2714", 10, GREEN, True)

slide_number(s, 9)
notes(s, "Three core concepts in one slide: tokens (the currency), parameters (the engine size), context (the memory). For cost decisions: more context = more tokens per request = higher cost.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — TEMPERATURE, RAG, PROMPTING
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Temperature, RAG & Prompt Engineering")

# Left: Temperature scale
add_textbox(s, M, Inches(0.9), Inches(4.0), Inches(0.4), "Temperature", 16, ACCENT, True)
add_textbox(s, M, Inches(1.35), Inches(4.0), Inches(0.3), "Controls creativity (0 = deterministic, 1+ = creative)", 10, MDGRAY)
temp_zones = [
    (Inches(1.0), Inches(0.7), "0.0", "Deterministic\nMath, code, facts", RGBColor(0x1A, 0x2D, 0x4D)),
    (Inches(2.2), Inches(1.4), "0.3\u20130.7", "Balanced\nMost tasks, docs, analysis", ACCENT),
    (Inches(3.4), Inches(1.4), "1.0+", "Creative\nBrainstorming, writing", GREEN),
]
for x, w, label, desc, color in temp_zones:
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), w, Inches(1.2), color, f"{label}\n\n{desc}", 9, WHITE)

# Center: RAG diagram
add_textbox(s, Inches(4.8), Inches(0.9), Inches(3.5), Inches(0.4), "What is RAG?", 16, ACCENT, True)
rag_y = Inches(1.5)
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), rag_y, Inches(1.5), Inches(0.7), ACCENT, "User asks\nquestion", 9, WHITE)
make_arrow(s, Inches(6.3), rag_y + Inches(0.35), Inches(7.0), rag_y + Inches(0.35), ACCENT)
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), rag_y, Inches(1.5), Inches(0.7), GREEN, "Search your\ndocuments", 9, WHITE)
make_arrow(s, Inches(7.0) + Inches(1.5), rag_y + Inches(0.35), Inches(7.0) + Inches(2.2), rag_y + Inches(0.35), ACCENT)
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), rag_y + Inches(1.0), Inches(1.5), Inches(0.7), RGBColor(0x16, 0x6B, 0x34), "LLM answers\nwith refs", 9, WHITE)
make_arrow(s, Inches(6.3), rag_y + Inches(1.35), Inches(7.0), rag_y + Inches(1.35), ACCENT)
add_textbox(s, Inches(4.8), Inches(3.9), Inches(3.8), Inches(0.8),
    "Think: \"Hey LLM, read these 3 docs\nand answer the question.\"\nFixes hallucinations.", 10, MDGRAY)

# Right: Prompt Engineering
add_textbox(s, Inches(9.2), Inches(0.9), Inches(3.5), Inches(0.4), "Prompt Engineering", 16, ACCENT, True)
add_bullets(s, Inches(9.2), Inches(1.35), Inches(3.5), Inches(2.0), [
    "Be specific",
    '  "Summarize in 3 bullet\n'
    '   points with action items"',
    "",
    "Give examples",
    '  Show 2-3 input/output\n'
    '   pairs first',
    "",
    "Ask for structure",
    '  "Output as JSON"\n'
    '  "Present as a table"',
], 10, BLACK, Pt(6), "\u25B6")

# Bottom: key insight
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.5), Inches(11.3), Inches(1.3), RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(s, Inches(1.3), Inches(5.6), Inches(10.7), Inches(1.1),
    'Temperature = creativity knob (low for code, high for brainstorming)  |  RAG = "read my docs before answering" — eliminates most hallucinations\n'
    'Prompt quality directly determines output quality. "Vague in, vague out." Being specific is the highest-ROI skill with AI.',
    11, GREEN)
slide_number(s, 10)
notes(s, "Three concepts that every AI user should understand. Temperature is why the same question gets different answers. RAG is the most practical way to make AI work with your company's data.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — 3 WAYS TO ADAPT
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "3 Ways to Adapt AI to Your Needs")
add_table(s, ["Method", "How It Works", "Time", "Cost", "Complexity", "Best For"],
    [["Prompting", "Write better instructions", "Instant", "Free", "Low", "Most tasks: writing, analysis, coding"],
     ["RAG", "Feed it your documents", "Hours\u2013Days", "Low", "Medium", "Docs Q&A, customer support, knowledge bases"],
     ["Fine-tuning", "Train it on your examples", "Days\u2013Weeks", "$100\u2013$10K", "High", "Specialized domains, brand tone, niche tasks"]],
    M, Inches(1.0), Inches(11.3), Inches(0.55))
# Guidance
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(3.3), Inches(11.3), Inches(0.9), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(3.4), Inches(10.7), Inches(0.7),
    'Start with prompting. Add RAG when you need the model to know YOUR data. Only fine-tune when you need specialized behavior at scale.',
    12, ACCENT, True)

# Use case cards
cases = [
    ("Writing emails & reports", "Prompting", GREEN),
    ("Customer support chatbot", "RAG", ACCENT),
    ("Legal document analysis", "Fine-tuning", RGBColor(0x1A, 0x2D, 0x4D)),
    ("Code generation", "Prompting", GREEN),
    ("Internal docs Q&A", "RAG", ACCENT),
    ("Brand voice content", "Fine-tuning", RGBColor(0x1A, 0x2D, 0x4D)),
]
for i, (task, method, color) in enumerate(cases):
    col = i % 3; row = i // 3
    x = M + col * Inches(4.0); y = Inches(4.6) + row * Inches(0.7)
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(0.55), WHITE, f"{task}  \u2192  {method}", 11, color)
    s.shapes[-1].line.color.rgb = LTGRAY; s.shapes[-1].line.width = Pt(1)

slide_number(s, 11)
notes(s, "Adaptation spectrum: prompting (free, instant) → RAG (add documents) → fine-tuning (train on examples). Most companies never need fine-tuning. Prompting + RAG covers 95% of use cases.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — MODEL ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Model Architecture & Scaling")
add_table(s, ["Architecture", "How It Works", "Examples", "Best For"],
    [["Encoder-only", "Processes input bidirectionally", "BERT, RoBERTa", "Classification, NER, search"],
     ["Decoder-only", "Generates autoregressively (left-to-right)", "GPT, Claude, Llama, Gemini", "Text generation, chat, coding"],
     ["Encoder-Decoder", "Encode input, decode output", "T5, BART, mT5", "Translation, summarization"],
     ["MoE (Mixture of Experts)", "8+ expert sub-networks, router picks top 2", "Llama 4, DeepSeek V4", "Scale params without inflating compute"]],
    M, Inches(1.0), Inches(11.3), Inches(0.55))
add_textbox(s, M, Inches(3.6), Inches(11.3), Inches(2.5),
    'All major modern LLMs are decoder-only.\n'
    'MoE = more parameters without proportionally more compute. Llama 4 has 405B total but only ~40B active per token.\n'
    'Training cost: Frontier models = $10M-$100M+. Fine-tuning existing models = $100-$10K.\n'
    'The trend: models getting smaller (edge/SLM), faster, cheaper — while matching frontier quality on specific tasks.',
    12, BLACK)
slide_number(s, 12)
notes(s, "Technical detail slide for the curious. Key point: decoder-only dominates modern LLMs. MoE is how they scale parameters without proportional cost.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — PART 3 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(3, "Tools & Practical Usage", "Chat, coding, content tools — and making build-vs-buy decisions")
slide_number(prs.slides[-1], 13)
notes(prs.slides[-1], "Now the practical part: which tools exist, what they cost, and how to choose the right approach for your organization.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — TOOL LANDSCAPE GRID
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "The AI Tool Landscape")
categories = [
    ("Chat Interfaces", ACCENT, ["Claude.ai — Best writing, reasoning", "ChatGPT — All-purpose, web search", "Gemini — 1M context, research", "Perplexity — Cited sources", "DeepSeek — Free/$10, reasoning", "Grok 3 — Real-time X data"]),
    ("Coding Tools", GREEN, ["Cursor — Best IDE ($20-40/mo)", "Claude Code — Best CLI (pay-per-token)", "GitHub Copilot — Best free tier", "Windsurf — Agentic workflows", "Aider — Git-native, open-source"]),
    ("Content Creation", AMBER, ["Midjourney — Photorealistic images", "Runway — Video editing/generation", "Suno — Music ($0-10/mo)", "ElevenLabs — Voice synthesis", "Sora — Cinematic video", "NotebookLM — Document analysis (free)"]),
    ("Infrastructure", RGBColor(0x1A, 0x2D, 0x4D), ["Anthropic / OpenAI / Google APIs", "Ollama + Llama (local, free)", "LangChain / CrewAI (open-source)", "Vector DBs: Pinecone, Chroma, pgvector", "Cloud: AWS Bedrock, Vertex AI, Azure AI"]),
]
for i, (cat_title, color, items) in enumerate(categories):
    col = i % 2; row = i // 2
    x = M + col * Inches(6.0); y = Inches(1.0) + row * Inches(3.0)
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(2.8), WHITE)
    s.shapes[-1].line.color.rgb = color; s.shapes[-1].line.width = Pt(2)
    add_textbox(s, x + Inches(0.2), y + Inches(0.1), Inches(5.0), Inches(0.35), cat_title, 14, color, True)
    add_textbox(s, x + Inches(0.2), y + Inches(0.5), Inches(5.0), Inches(2.2), '\n'.join(items), 10, BLACK)
slide_number(s, 14)
notes(s, "Four categories of AI tools. Most teams need one from each category: a chat tool, a coding tool, a content tool, and an API for building custom solutions.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 15 — BUILD VS BUY DECISION TREE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Build vs Buy Decision Framework")
# Decision tree using shapes
y = Inches(1.2)
# Start
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), y, Inches(4.3), Inches(0.7), ACCENT, "Q: Is data privacy critical\n(healthcare, legal, finance)?", 11, WHITE)
# Yes branch
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y + Inches(1.2), Inches(3.0), Inches(0.7), GREEN, "YES\n\nSelf-host open-weight model\n(Llama, Qwen, DeepSeek)\nFull data control", 10, WHITE)
make_arrow(s, Inches(4.5) + Inches(2.15), y + Inches(0.7), Inches(2.5), y + Inches(1.2), GREEN)
add_textbox(s, Inches(4.2), y + Inches(0.5), Inches(0.5), Inches(0.3), "Yes", 9, GREEN)
# No branch
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), y + Inches(1.2), Inches(3.0), Inches(0.7), ACCENT, "NO\n\nQ: Usage > 50M tokens/month?", 10, WHITE)
make_arrow(s, Inches(4.5) + Inches(4.3), y + Inches(0.35), Inches(6.5), y + Inches(1.2), ACCENT)
add_textbox(s, Inches(9.0), y + Inches(0.5), Inches(0.5), Inches(0.3), "No", 9, ACCENT)
# Self-host branch
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), y + Inches(2.4), Inches(3.0), Inches(0.7), GREEN, "YES\n\nSelf-host with auto-scaling\n(vLLM + 70B INT8)\nCost: $200-3K/mo", 10, WHITE)
make_arrow(s, Inches(5.5) + Inches(1.5), y + Inches(1.9), Inches(7.0), y + Inches(2.4), GREEN)
add_textbox(s, Inches(6.3), y + Inches(1.7), Inches(0.5), Inches(0.3), "Yes", 9, GREEN)
# API branch
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), y + Inches(2.4), Inches(3.0), Inches(0.7), ACCENT, "NO (< 50M)\n\nUse API (Claude/GPT/Gemini)\nCheapest & fastest option\nNo infra needed", 10, WHITE)
make_arrow(s, Inches(5.5) + Inches(3.0), y + Inches(1.55), Inches(10.2), y + Inches(2.4), ACCENT)
add_textbox(s, Inches(8.8), y + Inches(1.7), Inches(0.5), Inches(0.3), "No", 9, ACCENT)

# Bottom summary
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(4.0), Inches(11.3), Inches(2.8), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(4.2), Inches(10.7), Inches(2.4),
    'Rule of Thumb:\n'
    '• < 10M tokens/month: API — fastest, cheapest, zero infrastructure\n'
    '• 10\u201350M tokens/month: Evaluate API vs fine-tuned open model (break-even analysis)\n'
    '• > 50M tokens/month: Self-host saves significant money. Open-weight models (Llama, DeepSeek) are free\n'
    '• Data-sensitive: Self-host always. Privacy is the #1 reason companies self-host\n\n'
    'Example: 1M chatbot conversations/month with Claude API = ~$2.10 total. Engineering time matters more.',
    11, ACCENT)
slide_number(s, 15)
notes(s, "Decision tree for build vs buy. Most companies should use APIs. Only self-host when privacy demands it or volume exceeds 50M tokens/month.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 16 — API PRICING BAR CHART
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "API Pricing Comparison (May 2026)")

# Bar chart: input vs output pricing
chart_data = CategoryChartData()
chart_data.categories = ['Claude\nOpus 4.7', 'Claude\nSonnet', 'GPT-5.5', 'Gemini\n3.1 Pro', 'DeepSeek\nV4', 'Llama 4\n(self-host)']
chart_data.add_series('Input ($/1M tokens)', (15, 3, 2, 2, 0.14, 0))
chart_data.add_series('Output ($/1M tokens)', (75, 15, 8, 12, 0.28, 0))
chart = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, M, Inches(1.0), Inches(7.0), Inches(4.0), chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(9)
plot.data_labels.number_format = '$#,##0.00'
plot.data_labels.show_value = True
# Color the series
series_input = chart.series[0]; series_input.format.fill.solid(); series_input.format.fill.fore_color.rgb = ACCENT
series_output = chart.series[1]; series_output.format.fill.solid(); series_output.format.fill.fore_color.rgb = GREEN

# Right: key insights
add_textbox(s, Inches(8.5), Inches(1.0), Inches(4.0), Inches(4.5),
    'Key Pricing Insights\n\n'
    '\u25CF Claude Opus = best quality,\n'
    '   most expensive (50x DeepSeek)\n\n'
    '\u25CF GPT-5.5 & Gemini = great value\n'
    '   at $2/1M input tokens\n\n'
    '\u25CF DeepSeek V4 = 10-50x cheaper\n'
    '   than leading closed models\n\n'
    '\u25CF Llama 4 = free to self-host\n'
    '   (you pay for GPU servers)\n\n'
    '\u25CF For most apps, model cost is\n'
    '   NEGLIGIBLE vs engineering time',
    12, BLACK)
# Callout
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.4), Inches(11.3), Inches(1.3), RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(s, Inches(1.3), Inches(5.5), Inches(10.7), Inches(1.0),
    'Example: 100K input + 20K output with Claude Sonnet = $0.33/request. With DeepSeek = $0.014/request (23x cheaper).\n'
    '200K conversations/month with Sonnet = $66. Choose your model based on quality needs, not cost.',
    11, GREEN)
slide_number(s, 16)
notes(s, "Pricing bar chart. DeepSeek's cost advantage is staggering: $0.14 vs $15 per 1M input tokens. But price isn't everything — Claude and GPT offer better reliability and ecosystem.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 17 — TOOL PAIRING + BUILD
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Tool Pairing Strategies & What You Can Build")

# Left: Tool combos
add_textbox(s, M, Inches(0.9), Inches(5.5), Inches(0.4), "Recommended Tool Stack", 16, ACCENT, True)
combos = [
    ("Daily ($40-60/mo)", GREEN, ["Chat: Claude.ai + ChatGPT", "Code: Cursor + Claude Code", "Research: Gemini + Perplexity"]),
    ("Budget ($20-40/mo)", ACCENT, ["Free tiers of all major tools", "Add Cursor ($20/mo)", "Add one Pro plan ($20/mo)"]),
    ("Enterprise (per seat)", RGBColor(0x1A, 0x2D, 0x4D), ["Claude Team / ChatGPT Team", "Cursor Business / Copilot Enterprise", "Midjourney Pro + n8n workflow"]),
]
for i, (title, color, items) in enumerate(combos):
    y = Inches(1.6) + i * Inches(1.6)
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, y, Inches(5.5), Inches(1.3), WHITE, "", 10, BLACK)
    s.shapes[-1].line.color.rgb = color; s.shapes[-1].line.width = Pt(2)
    add_textbox(s, Inches(1.2), y + Inches(0.05), Inches(5.0), Inches(0.3), title, 13, color, True)
    add_textbox(s, Inches(1.2), y + Inches(0.4), Inches(5.0), Inches(0.8), '\n'.join(items), 10, BLACK)

# Right: What you can build
add_textbox(s, Inches(7.2), Inches(0.9), Inches(5.5), Inches(0.4), "What You Can Build Today", 16, ACCENT, True)
builds = [
    ("\U0001F4AC", "RAG Chatbot", "Q&A on your docs, knowledge base. Hours to build, deployed on day 1."),
    ("\U0001F4DD", "Content Pipeline", "Auto-generate blog posts, reports, summaries. 10x content output."),
    ("\U0001F50D", "Data Analysis", "Natural language SQL queries. Non-technical teams can explore data."),
    ("\U0001F4E6", "Customer Support", "Auto-answer 60-80% of tickets. Escalate complex ones to humans."),
    ("\U0001F4CA", "Research Assistant", "Summarize papers, market reports. Extract key findings instantly."),
    ("\u2699\uFE0F", "Code Automation", "CI/CD bots for code review, test generation, refactoring."),
]
for i, (icon, title, desc) in enumerate(builds):
    y = Inches(1.6) + i * Inches(0.75)
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), y, Inches(5.5), Inches(0.6), WHITE)
    s.shapes[-1].line.color.rgb = LTGRAY; s.shapes[-1].line.width = Pt(1)
    add_textbox(s, Inches(7.4), y + Inches(0.05), Inches(0.5), Inches(0.5), icon, 16, BLACK, False, PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.9), y + Inches(0.02), Inches(1.5), Inches(0.3), title, 11, ACCENT, True)
    add_textbox(s, Inches(7.9), y + Inches(0.3), Inches(4.6), Inches(0.3), desc, 9, MDGRAY)
slide_number(s, 17)
notes(s, "Tool pairing and use cases. Key message: you don't need everything. Start with free tiers, add tools as you need them. All 6 use cases achievable with existing 2026 tools.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 18 — GETTING STARTED ROADMAP
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Getting Started — Your First Month with AI")

# Roadmap timeline
steps = [
    ("30 Minutes", "Today", GREEN, ["Sign up at claude.ai or chatgpt.com (free)", "Ask: 'Explain X like I'm 12' then 'like a pro'", "Paste in an email → 'Summarize + suggest improvements'"]),
    ("1 Day", "This week", ACCENT, ["Pick your daily tool: Claude (writing) or ChatGPT (general)", "Try coding: 'Write a Python script to rename files'", "Explore: 3 prompts on a real task from your work"]),
    ("1 Week", "Next week", AMBER, ["Set up Cursor (IDE) + GitHub Copilot (free tier)", "Build a RAG prototype on your docs (LangChain/OpenAI)", "Create a custom GPT with your team's knowledge base"]),
    ("1 Month", "This quarter", RGBColor(0x1A, 0x2D, 0x4D), ["Deploy a production RAG chatbot or content pipeline", "Evaluate build-vs-buy for your highest-volume use case", "Train 2-3 team members on AI tool usage"]),
]
for i, (time, when, color, items) in enumerate(steps):
    x = M + i * Inches(3.0); y = Inches(1.3); w = Inches(2.7)
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.6), color, f"{time} \u2014 {when}", 12, WHITE)
    content = '\n'.join(['\u25CF ' + item for item in items])
    add_textbox(s, x, y + Inches(0.8), w, Inches(4.0), content, 10, BLACK)

# Bottom resources
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.5), Inches(11.3), Inches(1.3), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.6), Inches(10.7), Inches(1.1),
    'Free Resources: claude.ai  \u2022  chatgpt.com  \u2022  gemini.google.com  \u2022  cursor.com  \u2022  github.com/features/copilot  \u2022  ollama.com (local models)\n'
    'Learning: The AI Playbook (ai-playbook.pages.dev)  \u2022  fast.ai  \u2022  DeepLearning.AI  \u2022  3Blue1Brown (YouTube)',
    11, ACCENT)
slide_number(s, 18)
notes(s, "Concrete getting-started roadmap with specific actions at each stage. Make it feel achievable — 30 minutes today, gradual ramp over a month.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 19 — BEST PRACTICES
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Best Practices for Using AI")
# DO column
add_textbox(s, M, Inches(0.9), Inches(5.0), Inches(0.4), "DO", 18, GREEN, True)
dos = [
    "Verify important facts — LLMs can hallucinate",
    "Be specific in prompts — vague = vague answers",
    "Use AI as a starting point, not the final answer",
    "Experiment with different models for different tasks",
    "Keep humans in the loop for critical decisions",
    "Invest in prompt engineering skills — highest ROI",
]
for i, item in enumerate(dos):
    add_textbox(s, M, Inches(1.5) + i * Inches(0.45), Inches(5.0), Inches(0.4), f'\u2713  {item}', 11, GREEN)

# DON'T column
add_textbox(s, Inches(7.2), Inches(0.9), Inches(5.5), Inches(0.4), "DON'T", 18, RED, True)
donts = [
    "Paste sensitive/private data into public AI services",
    "Trust AI-generated code without review",
    "Assume the model knows recent events (check cutoff)",
    "Over-automate — some things need human judgment",
    "Use AI for legal/medical advice without expert review",
    "Forget: AI amplifies humans, doesn't replace them",
]
for i, item in enumerate(donts):
    add_textbox(s, Inches(7.2), Inches(1.5) + i * Inches(0.45), Inches(5.5), Inches(0.4), f'\u2715  {item}', 11, RED)

slide_number(s, 19)
notes(s, "Best practices distilled into actionable DOs and DON'Ts. The privacy point is critical — never paste confidential data into public AI services.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 20 — WHERE AI IS HEADING (2026-2028)
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Where AI is Heading — 2026 to 2028")
trends = [
    ("Agentic AI", "Models that plan, use tools, and execute multi-step tasks autonomously", ACCENT),
    ("Edge & SLMs", "Small models on phones/laptops. Zero-cost inference, zero latency", GREEN),
    ("Multimodal", "Vision, audio, video natively integrated. No more text-only", AMBER),
    ("Open Dominance", "Open models (Llama, DeepSeek, Qwen) matching closed models on all benchmarks", RGBColor(0x1A, 0x2D, 0x4D)),
    ("Reasoning Default", "Chain-of-thought, tree-of-thought become standard. Specialized reasoning models", RGBColor(0x93, 0x3D, 0xEA)),
    ("AI Ubiquity", "AI in email, docs, IDE, browser, terminal — embedded in every tool we use", RED),
]
for i, (title, desc, color) in enumerate(trends):
    col = i % 3; row = i // 3
    x = M + col * Inches(4.0); y = Inches(1.0) + row * Inches(2.8)
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.6), Inches(2.3), WHITE)
    s.shapes[-1].line.color.rgb = color; s.shapes[-1].line.width = Pt(2)
    add_textbox(s, x + Inches(0.2), y + Inches(0.1), Inches(3.2), Inches(0.4), title, 16, color, True)
    add_textbox(s, x + Inches(0.2), y + Inches(0.6), Inches(3.2), Inches(1.5), desc, 11, BLACK)
slide_number(s, 20)
notes(s, "Looking ahead 2-3 years. Key message: the pace is accelerating. What seems cutting-edge today will be table stakes in 18 months.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 21 — HUMAN + AI ADVANTAGE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "The Human + AI Advantage")

# AI column
add_textbox(s, M, Inches(1.0), Inches(5.0), Inches(0.4), "What AI Excels At", 16, ACCENT, True)
ai_skills = ["Speed — processes millions of tokens in seconds", "Scale — handles thousands of requests simultaneously",
    "Pattern Recognition — finds correlations humans miss", "Consistency — same quality 24/7, never tired",
    "Language — translates, summarizes, generates text effortlessly",
    "Learning from examples — adapts to new tasks with a prompt"]
add_bullets(s, M, Inches(1.5), Inches(5.0), Inches(3.0), ai_skills, 11, BLACK, Pt(10), "\u25CF")

# Human column
add_textbox(s, Inches(7.2), Inches(1.0), Inches(5.5), Inches(0.4), "What Humans Excel At", 16, GREEN, True)
human_skills = ["Judgment — knowing when a decision is right, not just logical", "Creativity — original ideas, not recombinations of existing ones",
    "Ethics — understanding context, nuance, and consequences", "Domain Expertise — deep specialized knowledge",
    "Empathy — understanding user needs, emotions, and motivations",
    "Strategy — setting direction, making tradeoffs, taking accountability"]
add_bullets(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.0), human_skills, 11, BLACK, Pt(10), "\u25CF")

# Center overlap: together wins
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.0), Inches(11.3), Inches(1.8), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.1), Inches(10.7), Inches(1.6),
    'The Winning Formula: AI drafts \u2192 human reviews  |  AI analyzes data \u2192 human makes decisions  |  AI generates ideas \u2192 human curates\n\n'
    'This is NOT about replacement. It\'s about amplification.\n'
    'Your job is not at risk from AI. Your job is at risk from someone who knows how to USE AI.',
    13, ACCENT, True)
slide_number(s, 21)
notes(s, "The closing argument. AI + humans outperform either alone. The key skill of the next decade: knowing how to leverage AI effectively.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 22 — KEY TAKEAWAYS
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(0.8), Inches(11.3), Inches(1.0), "Key Takeaways", 36, ACCENT, True, PP_ALIGN.LEFT)
takeaways = [
    ("1", "LLMs are pattern-matching at scale — not magic, just math trained on trillions of examples."),
    ("2", "Five key players: Anthropic, OpenAI, Google, DeepSeek, Meta. Each has distinct strengths."),
    ("3", "Prices dropping 50%+ per year. Yesterday's premium is today's free. DeepSeek is 100x cheaper than Claude."),
    ("4", "Better prompts = better outputs. Being specific is the single highest-ROI AI skill."),
    ("5", "Start with free tiers today. Claude.ai and ChatGPT are free. You can be productive in 30 minutes."),
    ("6", "RAG feeds your documents to AI — the most practical way to make AI useful for YOUR business."),
    ("7", "AI amplifies humans. It doesn't replace them. The winning formula: AI + human judgment."),
]
for i, (num, text) in enumerate(takeaways):
    y = Inches(1.8) + i * Inches(0.75)
    make_shape(s, MSO_SHAPE.OVAL, M, y, Inches(0.5), Inches(0.5), ACCENT, num, 14, WHITE)
    add_textbox(s, Inches(1.8), y + Inches(0.05), Inches(10.5), Inches(0.5), text, 16, BLACK)
slide_number(s, 22)
notes(s, "The 7 key takeaways. These should be memorable enough to repeat in the hallway after the presentation.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 23 — Q&A
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(2.0), Inches(11.3), Inches(1.5), "Questions?", 52, ACCENT, True, PP_ALIGN.CENTER)
add_textbox(s, M, Inches(3.5), Inches(11.3), Inches(1.0), "Ask me anything about AI, LLMs, or how to get started", 18, BLACK, False, PP_ALIGN.CENTER)
add_textbox(s, M, Inches(5.0), Inches(11.3), Inches(0.5), "ai-playbook.pages.dev  \u2022  claude.ai  \u2022  chatgpt.com  \u2022  gemini.google.com", 14, MDGRAY, False, PP_ALIGN.CENTER)
slide_number(s, 23)
notes(s, "Open Q&A. Be prepared for questions about: cost, security/privacy, when to use AI, specific tools.")

# ─── SAVE ───────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ai-playbook-intro.pptx")
prs.save(output_path)
print(f"\u2713 Saved {output_path}")
print(f"  {len(prs.slides)} slides  |  Corporate light theme  |  Charts + diagrams + tables")
