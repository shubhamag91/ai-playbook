#!/usr/bin/env python3
"""AI Playbook Intermediate Presentation Generator — Architecture & Systems Design."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

# Colors
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF8, 0xF9, 0xFA)  # Light warm gray bg
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
DKGRAY  = RGBColor(0x4A, 0x4A, 0x4A)
MDGRAY  = RGBColor(0x88, 0x88, 0x80)
LTGRAY  = RGBColor(0xE0, 0xDC, 0xD6)
ACCENT  = RGBColor(0x3B, 0x82, 0xF6)  # Blue
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
M = Inches(1.0)  # margin

TOTAL_SLIDES = 12

# ─── Helpers ───────────────────────────────────────────────
def new_slide(bg_color=BG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg_color
    return s

def header_bar(slide, text, color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = M
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return bar

def slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{num} / {TOTAL_SLIDES}"; p.font.size = Pt(9); p.font.color.rgb = MDGRAY; p.alignment = PP_ALIGN.RIGHT

def notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_textbox(slide, left, top, width, height, text, font_size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=14, color=BLACK, spacing=Pt(12), bullet_char="\u25CF"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size); p.font.color.rgb = color
        p.space_after = spacing
    return txBox

def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.45), font_size=11):
    n_rows = len(rows) + 1; n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows).table
    col_w = (int(width) - 300000) / n_cols
    for i in range(n_cols): tbl.columns[i].width = int(col_w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs: p.font.size = Pt(font_size-1); p.font.color.rgb = WHITE; p.font.bold = True
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(val)
            for p in c.text_frame.paragraphs: p.font.size = Pt(font_size); p.font.color.rgb = BLACK
            if i % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = LTGRAY
    return tbl

def make_shape(slide, shape_type, left, top, width, height, fill_color=ACCENT, text="", font_size=10, font_color=WHITE):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    if text:
        tf = shp.text_frame; tf.word_wrap = True; tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(font_size); p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
    return shp

def make_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color; connector.line.width = Pt(2)
    return connector

def metric_card(slide, left, top, width, height, number, label, color=ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LTGRAY; card.line.width = Pt(1)
    tf = card.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.text = str(number); p1.font.size = Pt(24); p1.font.color.rgb = color
    p1.font.bold = True; p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(10); p2.font.color.rgb = MDGRAY
    p2.alignment = PP_ALIGN.CENTER
    return card

def divider_slide(section_num, title, subtitle=""):
    s = new_slide(ACCENT)
    add_textbox(s, M, Inches(2.0), Inches(11.3), Inches(1.5), f"Part {section_num}", 18, WHITE, False, PP_ALIGN.LEFT)
    add_textbox(s, M, Inches(3.0), Inches(11.3), Inches(1.5), title, 36, WHITE, True, PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(s, M, Inches(4.2), Inches(11.3), Inches(1.0), subtitle, 16, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.LEFT)
    return s

# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(1.8), Inches(11.3), Inches(1.5), "AI System Design: Building with LLMs", 44, ACCENT, True, PP_ALIGN.LEFT)
add_textbox(s, M, Inches(3.2), Inches(11.3), Inches(1.0), "From Prompting to RAG, Fine-tuning, & Infrastructure Decisions", 22, BLACK, False, PP_ALIGN.LEFT)
add_textbox(s, M, Inches(4.2), Inches(11.3), Inches(0.6), "Level 2: Builder Presentation  \u2022  AI Playbook", 13, MDGRAY, False, PP_ALIGN.LEFT)
make_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(3.1), Inches(1.8), Pt(4), ACCENT)
slide_number(s, 1)
notes(s, "Welcome to the Builder Level Presentation. Today we will dive under the hood of LLM application design, review the training pipeline, contrast adaptation styles, inspect RAG architectures, and map build-vs-buy parameters.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — PART 1 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(1, "Deepening the Mechanics", "Tokens, training pipelines, and the LLM adaptation spectrum")
slide_number(prs.slides[-1], 2)
notes(prs.slides[-1], "In this section, we move beyond simple autocomplete to inspect tokens, parameters, and the formal training pipelines.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — TOKENS, CONTEXT & MEMORY
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Tokens, Context Windows & Memory limits")

col_w = Inches(5.3); col_x = [M, Inches(7.0)]

# Left: Tokens
add_textbox(s, col_x[0], Inches(1.0), col_w, Inches(0.4), "Tokens: The Currency of AI", 18, ACCENT, True)
add_textbox(s, col_x[0], Inches(1.5), col_w, Inches(3.0),
    "Models do not process characters or words directly; they process tokens.\n\n"
    "\u2022 A token is approximately \u00BE of a word.\n"
    "\u2022 1,000 tokens \u2248 750 words \u2248 1.5 single-spaced pages.\n"
    "\u2022 Input Tokens: Prompt data read by the model (cheaper).\n"
    "\u2022 Output Tokens: Response data written by the model (more expensive).\n\n"
    "Example text segmentation:",
    12, BLACK)
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, col_x[0], Inches(4.7), col_w, Inches(0.7), LTGRAY)
add_textbox(s, col_x[0], Inches(4.8), col_w, Inches(0.5), "Token | i | za | tion |  helps |  mod | els |  read", 12, BLACK, False, PP_ALIGN.CENTER)

# Right: Context
add_textbox(s, col_x[1], Inches(1.0), col_w, Inches(0.4), "Context Windows: Scaling Conversation Memory", 18, ACCENT, True)
add_textbox(s, col_x[1], Inches(1.5), col_w, Inches(3.0),
    "The context window determines the exact amount of text the model can process in a single conversation history.\n\n"
    "Crucially: The entire chat log is re-sent with every new message. Longer histories cost exponentially more in token fees.\n\n"
    "2026 Memory Windows compared:\n"
    "\u2022 GPT-5.5: 128,000 tokens (~200 pages)\n"
    "\u2022 Claude Opus 4.7: 400,000 tokens (~600 pages)\n"
    "\u2022 Gemini 3.1 Pro: 1,000,000 tokens (~1,500 pages / 2x War & Peace)",
    12, BLACK)
add_textbox(s, col_x[1], Inches(4.7), col_w, Inches(0.5), "Gemini (1M) > Claude (400K) > GPT (128K)", 11, GREEN, True, PP_ALIGN.CENTER)

slide_number(s, 3)
notes(s, "Explain that tokens are sub-word units. Large context windows mean models can analyze entire books or large repositories of code. However, longer history means higher per-message API costs since the context gets re-processed.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW LLMs LEARN (TRAINING PIPELINE)
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "How LLMs Learn — Deployed Training Pipeline")

# Visual representation of pipeline
stages = [
    ("Pretraining", "Trillions of words", "Unsupervised", "Cost: $10M-$100M", "Goal: Learn grammar, syntax, world facts"),
    ("SFT Fine-Tuning", "Thousands of target examples", "Supervised", "Cost: $100-$10K", "Goal: Adapt to custom tasks/voices"),
    ("RLHF / DPO", "Human preference pairings", "Reinforcement", "Cost: $1K-$50K", "Goal: Align outputs for helpfulness/safety")
]

for i, (title, data, mode, cost, goal) in enumerate(stages):
    x = M + i * Inches(3.9); y = Inches(1.2)
    card = make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(3.6), WHITE)
    card.line.color.rgb = ACCENT if i != 2 else GREEN; card.line.width = Pt(2)
    
    add_textbox(s, x + Inches(0.2), y + Inches(0.2), Inches(3.1), Inches(0.45), title, 16, ACCENT if i != 2 else GREEN, True)
    add_textbox(s, x + Inches(0.2), y + Inches(0.7), Inches(3.1), Inches(0.3), f"Mode: {mode}", 11, MDGRAY)
    
    # Details card
    make_shape(s, MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(1.1), Inches(3.1), Pt(1), LTGRAY)
    
    add_textbox(s, x + Inches(0.2), y + Inches(1.3), Inches(3.1), Inches(2.0),
        f"\u2022 Data size: {data}\n\n"
        f"\u2022 Training cost: {cost}\n\n"
        f"\u2022 Primary goal: {goal}", 10.5, BLACK)
    
    # Arrow between stages
    if i < 2:
        make_arrow(s, x + Inches(3.5), y + Inches(1.8), x + Inches(3.9), y + Inches(1.8), ACCENT)

# takeaway callout
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.1), Inches(11.3), Inches(1.3), RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(s, Inches(1.3), Inches(5.2), Inches(10.7), Inches(1.1),
    "Strategic Insight: You rarely need to pretrain a model. Fine-tuning is 10,000x cheaper. SFT teaches the model formatting and style, while RLHF (or DPO—Direct Preference Optimization) aligns it to human guidelines.",
    11, GREEN, True)

slide_number(s, 4)
notes(s, "Explain the training hierarchy. Pretraining is building the brain. SFT is teaching it a specific profession. RLHF is teaching it how to behave politely in office environments.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — ADAPTABILITY MATRIX
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Adaptability: Prompting vs. RAG vs. Fine-Tuning")

# Custom Table
headers = ["Method", "Set-up Time", "Compute Cost", "Setup Difficulty", "Data Access", "Best Suited For"]
rows = [
    ["Prompt Engineering", "Instant (minutes)", "Zero extra cost", "Low (writing skills)", "Static / In-context", "General writing, summaries, coding helper"],
    ["Retrieval-Augmented (RAG)", "Hours - Days", "Low ($10-$100/mo)", "Medium (vector APIs)", "Dynamic / Live files", "Customer support bots, Q&A on private files"],
    ["Fine-Tuning (SFT)", "Days - Weeks", "$100 - $10,000", "High (GPU training)", "Frozen in weights", "Highly specialized vocabulary, custom brand tone"]
]
tbl = add_table(s, headers, rows, M, Inches(1.1), Inches(11.3), Inches(0.55), 10.5)

# Highlight cells or add context
add_textbox(s, M, Inches(3.2), Inches(11.3), Inches(2.2),
    "The Golden Rule of AI Engineering:\n"
    "1. Earn the right to complexity. Start by writing rich, multi-shot prompt instructions.\n"
    "2. If the model lacks access to your corporate documents or live databases, build a RAG pipeline.\n"
    "3. Only attempt fine-tuning if the model fails to capture specialized formatting, vocabularies, or strict brand tone.",
    12, BLACK)

make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(0.8), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.95), Inches(10.7), Inches(0.5),
    "Recommendation: 95% of corporate application needs can be solved perfectly by pairing Prompting + RAG.", 11, ACCENT, True)

slide_number(s, 5)
notes(s, "The adaptation spectrum. Remind builders that fine-tuning does not prevent hallucinations, it only changes the model's formatting style. RAG is the primary cure for factual hallucination.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — PART 2 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(2, "RAG & Decision Frameworks", "Understanding retrieval-augmented architectures and build-vs-buy parameters")
slide_number(prs.slides[-1], 6)
notes(prs.slides[-1], "In the second part, we explore the mechanics of RAG (Retrieval-Augmented Generation) and analyze corporate system architectures.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — RAG LIFECYCLE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "What is RAG? The Retrieval-Augmented Generation Architecture")

# RAG Process diagram blocks
y = Inches(1.8); w_box = Inches(2.1); h_box = Inches(1.5); gap = Inches(0.18)

blocks = [
    ("1. User Query", "User inputs question in natural language.", RGBColor(0x1A, 0x2D, 0x4D)),
    ("2. Search Index", "System queries a Vector Database for context.", ACCENT),
    ("3. Retrieval", "Top matching text snippets are retrieved.", RGBColor(0x60, 0xA5, 0xFA)),
    ("4. Prompt Context", "Query + retrieved context are fused together.", GREEN),
    ("5. LLM Answer", "LLM responds using only the provided facts.", RGBColor(0x16, 0x6B, 0x34))
]

for i, (label, desc, color) in enumerate(blocks):
    x = M + i * (w_box + gap + Inches(0.1))
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w_box, h_box, color, label, 12, WHITE)
    
    # Description box below
    add_textbox(s, x, y + Inches(1.65), w_box, Inches(0.85), desc, 9.5, DKGRAY, False, PP_ALIGN.CENTER)
    
    # Connector arrows
    if i < len(blocks) - 1:
        make_arrow(s, x + w_box, y + h_box/2, x + w_box + gap + Inches(0.1), y + h_box/2, ACCENT)

# bottom box
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.1), Inches(11.3), Inches(1.3), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.2), Inches(10.7), Inches(1.1),
    "Analogy: RAG operates like an open-book exam. Instead of relying on what the model learned during training, the system searches the textbook for the exact page, feeds it to the model, and asks it to summarize the answer.",
    11, ACCENT)

slide_number(s, 7)
notes(s, "Explain why RAG is popular: it provides a clear audit trail. Because the model must cite its source snippet from the vector database, we can immediately verify its statements and prevent hallucinated fabrications.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEV & TOOL LANDSCAPE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "The AI Developer & Tool Landscape")

# 4 Quadrants
quads = [
    ("Chat Interfaces (SaaS)", ACCENT, 
     ["Claude.ai — Best text and formatting", 
      "ChatGPT Plus — Speed and web searches", 
      "Perplexity — Search index citation focus",
      "NotebookLM — Auto audio & summary notes"]),
    
    ("Coding Assistants (IDEs)", GREEN,
     ["Cursor — Market-leading AI IDE code editor", 
      "Claude Code — Terminal CLI-based automation", 
      "GitHub Copilot — Highly capable basic auto-complete",
      "Windsurf — Agentic multi-file code editing"]),
      
    ("Content Creation", AMBER,
     ["Midjourney — Best image synthesis", 
      "Runway / Sora — Professional-grade video generation",
      "Suno / Udio — AI audio generation", 
      "ElevenLabs — Hyper-realistic voice synthesis"]),
      
    ("Infrastructure & APIs", RGBColor(0x1A, 0x2D, 0x4D),
     ["APIs: Anthropic, OpenAI, Vertex AI", 
      "Local Hosting: Ollama, vLLM", 
      "Vector DBs: Pinecone, pgvector, Qdrant",
      "Frameworks: LangChain, LlamaIndex, CrewAI"])
]

for i, (title, color, items) in enumerate(quads):
    col = i % 2; row = i // 2
    x = M + col * Inches(6.0); y = Inches(1.0) + row * Inches(3.0)
    
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(2.8), WHITE)
    s.shapes[-1].line.color.rgb = color; s.shapes[-1].line.width = Pt(2)
    
    add_textbox(s, x + Inches(0.2), y + Inches(0.15), Inches(5.1), Inches(0.4), title, 14, color, True)
    add_bullets(s, x + Inches(0.2), y + Inches(0.65), Inches(5.1), Inches(2.0), items, 10.5, BLACK, Pt(4), "\u25AA")

slide_number(s, 8)
notes(s, "Categorize the tools for the builder. Remind them that most organizations combine multiple: Cursor for developers, Claude for corporate writers, and APIs + Vector Databases for custom internal software.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — BUILD VS BUY DECISION FRAMEWORK
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Build vs. Buy / Self-Hosting vs. API Decision Tree")

# Custom styled decision tree
y = Inches(1.1)
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), y, Inches(4.3), Inches(0.65), ACCENT, "Does data privacy or regulations\nstrictly forbid cloud hosting?", 11, WHITE)

# Branches
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y + Inches(1.0), Inches(3.0), Inches(0.7), GREEN, "YES\n\nSelf-Host Open-Weight Models\n(Llama 4, DeepSeek V4) on private VPC", 9.5, WHITE)
make_arrow(s, Inches(4.5) + Inches(2.15), y + Inches(0.65), Inches(2.5), y + Inches(1.0), GREEN)

make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), y + Inches(1.0), Inches(3.0), Inches(0.7), ACCENT, "NO\n\nIs monthly volume > 50M tokens?", 10, WHITE)
make_arrow(s, Inches(4.5) + Inches(4.3), y + Inches(0.35), Inches(6.5), y + Inches(1.0), ACCENT)

# Yes self-host
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), y + Inches(2.1), Inches(3.0), Inches(0.7), GREEN, "YES\n\nSelf-Host to reduce token fees\n(vLLM + Llama 70B on private server)", 9.5, WHITE)
make_arrow(s, Inches(5.5) + Inches(1.5), y + Inches(1.7), Inches(7.0), y + Inches(2.1), GREEN)

# No API
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), y + Inches(2.1), Inches(3.0), Inches(0.7), ACCENT, "NO\n\nUse Cloud APIs (Anthropic / OpenAI)\nFastest setup, zero infra maintenance", 9.5, WHITE)
make_arrow(s, Inches(5.5) + Inches(3.0), y + Inches(1.35), Inches(10.2), y + Inches(2.1), ACCENT)

# Decision breakdown table
headers_tree = ["Scenario", "Privacy Policy", "Maintenance Effort", "Estimated Monthly Cost", "Recommended Path"]
rows_tree = [
    ["Low volume (<10M tokens)", "Cloud permitted", "Zero", "Under $50 / month", "Use Cloud APIs (Sonnet/GPT)"],
    ["High volume (>50M tokens)", "Cloud permitted", "Medium (Docker setups)", "Cost break-even favors server", "vLLM / DeepSeek hosting"],
    ["Regulated (Medical/Legal)", "Local network only", "High (requires GPU admins)", "Hardware GPU lease fees", "Ollama / Llama on-premise"]
]
add_table(s, headers_tree, rows_tree, M, Inches(4.0), Inches(11.3), Inches(0.48), 10)

slide_number(s, 9)
notes(s, "Explain the decision mechanics. Unless highly restricted by medical or financial compliance, starting with Cloud APIs is almost always the right financial choice. Self-hosting only becomes cheaper at extreme token volumes.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — API PRICING BAR CHART
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "2026 API Pricing Comparison: Input vs. Output")

# Chart Data
chart_data = CategoryChartData()
chart_data.categories = ['Claude Opus 4.7', 'Claude Sonnet', 'GPT-5.5', 'Gemini Pro', 'DeepSeek V4', 'Llama 4 (on GPU)']
chart_data.add_series('Input ($/1M tokens)', (15, 3, 2, 2, 0.14, 0))
chart_data.add_series('Output ($/1M tokens)', (75, 15, 8, 12, 0.28, 0))

chart_left = M
chart_top = Inches(1.1)
chart_width = Inches(7.2)
chart_height = Inches(4.5)

chart = s.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, chart_left, chart_top, chart_width, chart_height, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# Colors
series_input = chart.series[0]; series_input.format.fill.solid(); series_input.format.fill.fore_color.rgb = ACCENT
series_output = chart.series[1]; series_output.format.fill.solid(); series_output.format.fill.fore_color.rgb = GREEN

# Right side: Pricing analysis
add_textbox(s, Inches(8.5), Inches(1.1), Inches(4.0), Inches(4.5),
    "Cost Dynamics Analysis:\n\n"
    "\u25AA Claude Opus 4.7 is the gold standard but costs 100x more than DeepSeek.\n\n"
    "\u25AA Sonnet and GPT-5.5 provide excellent builder value at $2-$3 per million tokens.\n\n"
    "\u25AA DeepSeek V4 matches performance on major benchmarks while costing a tiny fraction ($0.14).\n\n"
    "\u25AA Llama 4 is free to license, but you must factor in the cost of GPU server maintenance.",
    11.5, BLACK)

# Callout at bottom
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(0.8), RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(s, Inches(1.3), Inches(5.95), Inches(10.7), Inches(0.55),
    "Takeaway: Output tokens are typically 3x-5x more expensive than input tokens. Size your generations carefully.", 11, GREEN, True)

slide_number(s, 10)
notes(s, "Explain the pricing chart. Highlight that DeepSeek has completely disrupted the market by offering frontier capabilities at 100x lower fees. Llama has zero license fee but requires renting hardware GPUs.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — PRACTICAL BUILDER ROADMAP
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Practical Roadmap: Deploying Your First AI Pipeline")

phases = [
    ("Phase 1: Prototyping", "Duration: Days 1 - 5", 
     "\u2022 Setup API keys for Sonnet/GPT-5.5\n"
     "\u2022 Write markdown prompt templates\n"
     "\u2022 Use Cursor to draft code helpers\n"
     "\u2022 Verify basic output quality by hand", 
     ACCENT),
    
    ("Phase 2: RAG Pipeline", "Duration: Days 6 - 15",
     "\u2022 Connect pgvector or Pinecone DB\n"
     "\u2022 Write text chunking scripts\n"
     "\u2022 Build a semantic lookup system\n"
     "\u2022 Link lookup context to LLM prompts",
     GREEN),
     
    ("Phase 3: Production", "Duration: Days 16 - 30",
     "\u2022 Deploy API in a Docker container\n"
     "\u2022 Set strict token rate-limits\n"
     "\u2022 Add logging dashboards (logs.js)\n"
     "\u2022 Launch beta version for user testing",
     AMBER)
]

for i, (title, duration, bullets, color) in enumerate(phases):
    x = M + i * Inches(3.9); y = Inches(1.2)
    card = make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(4.3), WHITE)
    card.line.color.rgb = color; card.line.width = Pt(2)
    
    add_textbox(s, x + Inches(0.2), y + Inches(0.2), Inches(3.1), Inches(0.4), title, 15, color, True)
    add_textbox(s, x + Inches(0.2), y + Inches(0.55), Inches(3.1), Inches(0.25), duration, 10, MDGRAY)
    
    make_shape(s, MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.9), Inches(3.1), Pt(1), LTGRAY)
    
    add_textbox(s, x + Inches(0.2), y + Inches(1.1), Inches(3.1), Inches(3.0), bullets, 11, BLACK)

slide_number(s, 11)
notes(s, "Walk builders through this practical 30-day plan. Building doesn't require months. By breaking the task into Prototyping, RAG context injection, and Production telemetry, a solid MVP can be shipped in under a month.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION & Q&A
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(1.8), Inches(11.3), Inches(1.0), "Systems Design Q&A", 52, ACCENT, True, PP_ALIGN.CENTER)
add_textbox(s, M, Inches(3.0), Inches(11.3), Inches(0.8), "Ask me anything about RAG architectures, model selections, or API costs", 18, BLACK, False, PP_ALIGN.CENTER)
add_textbox(s, M, Inches(4.5), Inches(11.3), Inches(0.6), "Visit the Playbook: ai-playbook.pages.dev  \u2022  Cursor  \u2022  Pinecone", 14, MDGRAY, False, PP_ALIGN.CENTER)
slide_number(s, 12)
notes(s, "Open floor for builder-focused Q&A. Be prepared to address questions regarding token limits, vector database pricing, and handling raw document formatting inside RAG loaders.")

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ai-playbook-intermediate.pptx")
prs.save(output_path)
print(f"\u2713 Saved {output_path}")
