#!/usr/bin/env python3
"""AI Playbook Beginner Presentation Generator — Foundations & Basic Prompting."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
import os

# Colors
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF8, 0xF9, 0xFA)  # Light warm gray bg
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
DKGRAY  = RGBColor(0x4A, 0x4A, 0x4A)
MDGRAY  = RGBColor(0x88, 0x88, 0x80)
LTGRAY  = RGBColor(0xE0, 0xDC, 0xD6)
ACCENT  = RGBColor(0x3B, 0x82, 0xF6)  # Warm corporate blue
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
M = Inches(1.0)  # margin

TOTAL_SLIDES = 10

# ─── Helpers ───────────────────────────────────────────────
def new_slide(bg_color=BG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg_color
    return s

def header_bar(slide, text, color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = M
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return bar

def slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{num} / {TOTAL_SLIDES}"; p.font.size = Pt(9); p.font.color.rgb = MDGRAY; p.alignment = PP_ALIGN.RIGHT

def notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_textbox(slide, left, top, width, height, text, font_size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=14, color=BLACK, spacing=Pt(12), bullet_char="\u25CF"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size); p.font.color.rgb = color
        p.space_after = spacing
    return txBox

def make_shape(slide, shape_type, left, top, width, height, fill_color=ACCENT, text="", font_size=10, font_color=WHITE):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    if text:
        tf = shp.text_frame; tf.word_wrap = True; tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(font_size); p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
    return shp

def make_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color; connector.line.width = Pt(2)
    return connector

def metric_card(slide, left, top, width, height, number, label, color=ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LTGRAY; card.line.width = Pt(1)
    tf = card.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.text = str(number); p1.font.size = Pt(24); p1.font.color.rgb = color
    p1.font.bold = True; p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(10); p2.font.color.rgb = MDGRAY
    p2.alignment = PP_ALIGN.CENTER
    return card

def divider_slide(section_num, title, subtitle=""):
    s = new_slide(ACCENT)
    add_textbox(s, M, Inches(2.0), Inches(11.3), Inches(1.5), f"Part {section_num}", 18, WHITE, False, PP_ALIGN.LEFT)
    add_textbox(s, M, Inches(3.0), Inches(11.3), Inches(1.5), title, 36, WHITE, True, PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(s, M, Inches(4.2), Inches(11.3), Inches(1.0), subtitle, 16, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.LEFT)
    return s

# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(1.8), Inches(11.3), Inches(1.5), "AI & LLMs: The Basics for Everyone", 44, ACCENT, True, PP_ALIGN.LEFT)
add_textbox(s, M, Inches(3.2), Inches(11.3), Inches(1.0), "A Non-Technical Grounding in Generative AI", 22, BLACK, False, PP_ALIGN.LEFT)
add_textbox(s, M, Inches(4.2), Inches(11.3), Inches(0.6), "Level 1: Beginner Presentation  \u2022  AI Playbook", 13, MDGRAY, False, PP_ALIGN.LEFT)
make_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(3.1), Inches(1.8), Pt(4), ACCENT)
slide_number(s, 1)
notes(s, "Welcome to the Beginner Presentation. Today, we will unpack what Artificial Intelligence and Large Language Models actually are, debunk common myths, review the modern landscape, and cover practical methods to get started.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — PART 1 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(1, "The AI Landscape", "De-mystifying generative AI and explaining what is real vs. hype")
slide_number(prs.slides[-1], 2)
notes(prs.slides[-1], "In this first part, we will explore the AI landscape, clarify core concepts, and build a solid conceptual foundation.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS AN LLM?
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "What is an LLM? Autocomplete on Steroids")
add_textbox(s, M, Inches(1.0), Inches(5.2), Inches(5.0),
    "A Large Language Model (LLM) is software that predicts the next word in a sentence.\n\n"
    "It acts like the autocomplete feature on your mobile phone, but trained on massive amounts of data from the internet, books, and code.\n\n"
    "Key mechanics of the scale:\n"
    "\u2022 Trillions of words of training data\n"
    "\u2022 Billions of prediction parameters ('dials')\n"
    "\u2022 Massive memory windows (context)",
    13, BLACK)

# Autocomplete visual
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.5), Inches(1.2), WHITE)
s.shapes[-1].line.color.rgb = LTGRAY; s.shapes[-1].line.width = Pt(1)
add_textbox(s, Inches(7.0), Inches(1.3), Inches(5.1), Inches(0.4), "The best thing about an AI assistant is its ability to...", 14, DKGRAY)
add_textbox(s, Inches(7.0), Inches(1.8), Inches(5.1), Inches(0.4), "[ help ]   [ listen ]   [ write ]", 14, ACCENT, True)

# Core dimensions
dims = [
    ("Parameters", "Dials that control pattern matches", "10B - 1.7T"),
    ("Context", "Conversation memory window", "128K - 1M tokens"),
    ("Tokens", "Basic units of text processing", "1 token \u2248 \u00BE word")
]
for i, (title, desc, stat) in enumerate(dims):
    y = Inches(3.0) + i * Inches(1.2)
    card = make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y, Inches(5.5), Inches(1.0), WHITE)
    card.line.color.rgb = ACCENT; card.line.width = Pt(1.5)
    add_textbox(s, Inches(7.0), y + Inches(0.1), Inches(3.2), Inches(0.3), title, 13, ACCENT, True)
    add_textbox(s, Inches(7.0), y + Inches(0.45), Inches(3.2), Inches(0.4), desc, 10, MDGRAY)
    add_textbox(s, Inches(10.2), y + Inches(0.2), Inches(1.9), Inches(0.5), stat, 18, BLACK, True, PP_ALIGN.RIGHT)

slide_number(s, 3)
notes(s, "Think of LLMs as highly advanced autocomplete engines. By recognizing patterns across trillions of words, they generate text. They do not have feelings or sentience—they are math and pattern-matching at a massive scale.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — REALITY VS HYPE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Reality Check: Myths vs. Actual Capabilities")

# Left Column: Myths
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(1.0), Inches(5.3), Inches(4.8), WHITE)
s.shapes[-1].line.color.rgb = RED; s.shapes[-1].line.width = Pt(2)
add_textbox(s, Inches(1.2), Inches(1.2), Inches(4.4), Inches(0.45), "Common Myths (\u2715)", 16, RED, True)
myths = [
    "LLMs understand language like humans do.\n(They recognize and reproduce patterns.)",
    "ChatGPT is always accurate and factual.\n(They hallucinate and make up plausible facts.)",
    "More parameters automatically means smarter.\n(Data quality and architecture matter more.)",
    "AI assistants will immediately replace humans.\n(They augment humans; human oversight is vital.)"
]
add_bullets(s, Inches(1.2), Inches(1.8), Inches(4.4), Inches(4.2), myths, 11, BLACK, Pt(14), "\u2715")

# Right Column: Realities
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.0), Inches(5.3), Inches(4.8), WHITE)
s.shapes[-1].line.color.rgb = GREEN; s.shapes[-1].line.width = Pt(2)
add_textbox(s, Inches(7.2), Inches(1.2), Inches(4.9), Inches(0.45), "Actual Capabilities (\u2713)", 16, GREEN, True)
realities = [
    "Drafting, editing, and restructuring copy.",
    "Explaining complex subjects at various user levels.",
    "Generating and debugging programming code.",
    "Summarizing hundreds of pages into bullet points.",
    "Translating between natural languages and code.",
    "Tireless brainstorming partner for ideation."
]
add_bullets(s, Inches(7.2), Inches(1.8), Inches(4.9), Inches(4.2), realities, 11, BLACK, Pt(12), "\u2713")

# Callout at bottom
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(6.0), Inches(11.3), Inches(0.8), RGBColor(0xFE, 0xF2, 0xF2))
add_textbox(s, Inches(1.3), Inches(6.1), Inches(10.7), Inches(0.6),
    "Rule of Thumb: Never trust critical answers from an LLM blindly. Always verify facts.", 11, RED, True)

slide_number(s, 4)
notes(s, "Setting realistic expectations is crucial. LLMs are narrow assistants. They are world-class at generating, editing, and summarizing text, but they can and will hallucinate occasionally. Verify anything critical.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE 2026 AI LANDSCAPE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "The 2026 AI Landscape: Key Players & Costs")

# Metric cards
metric_card(s, M, Inches(1.0), Inches(2.6), Inches(1.0), "50%+", "Price Drop per Year", GREEN)
metric_card(s, Inches(3.8), Inches(1.0), Inches(2.6), Inches(1.0), "1 Million", "Max Memory (Tokens)", ACCENT)
metric_card(s, Inches(6.6), Inches(1.0), Inches(2.6), Inches(1.0), "87%", "Devs Using AI Tools", AMBER)
metric_card(s, Inches(9.4), Inches(1.0), Inches(2.9), Inches(1.0), "$0.14/1M", "Cheapest Model cost", GREEN)

# Core players
add_textbox(s, M, Inches(2.2), Inches(11.3), Inches(0.4), "Major AI Assistants & Services", 16, ACCENT, True)
players = [
    ("Claude (Anthropic)", "Outstanding at writing, logical reasoning, and complex coding. High quality."),
    ("ChatGPT (OpenAI)", "Excellent general-purpose utility, very fast, integrated web search features."),
    ("Gemini (Google)", "Largest memory window (up to 1M tokens), perfect for summarizing entire books."),
    ("DeepSeek V4", "Astonishingly inexpensive open-weight model matching closed systems."),
    ("Llama (Meta)", "Open-source foundation, free to self-host and customize completely.")
]
for i, (name, desc) in enumerate(players):
    x = M + (i % 2) * Inches(5.8); y = Inches(2.8) + (i // 2) * Inches(1.0)
    card = make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(0.85), WHITE)
    card.line.color.rgb = LTGRAY; card.line.width = Pt(1)
    add_textbox(s, x + Inches(0.2), y + Inches(0.08), Inches(5.1), Inches(0.25), name, 12, ACCENT, True)
    add_textbox(s, x + Inches(0.2), y + Inches(0.35), Inches(5.1), Inches(0.45), desc, 9.5, DKGRAY)

# bottom note
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.0), Inches(11.3), Inches(0.8), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(s, Inches(1.3), Inches(5.1), Inches(10.7), Inches(0.6),
    "Note: AI costs are declining incredibly fast, and open-weight models are making high-quality AI universally accessible.", 11, ACCENT)

slide_number(s, 5)
notes(s, "In 2026, AI is a commodity. Prices drop over 50% yearly, memory sizes have ballooned to 1M tokens, and open models like Llama or DeepSeek offer state-of-the-art capability for free or pennies.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — PART 2 DIVIDER
# ══════════════════════════════════════════════════════════════════════
divider_slide(2, "Prompting & Basic Workflows", "How to communicate with language models effectively for daily work")
slide_number(prs.slides[-1], 6)
notes(prs.slides[-1], "In Part 2, we will look at prompt engineering—the art and science of writing high-quality instructions to get high-quality answers.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — PROMPTING PRINCIPLES
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Three Core Prompting Principles")

principles = [
    ("1. Be Highly Specific", 
     "Tell the model exactly what you want.\n\n"
     "Bad: 'Write a summary of this.'\n"
     "Good: 'Write a 3-bullet-point summary focusing only on the financial action items.'", 
     ACCENT),
    ("2. Provide Examples (Few-Shot)", 
     "Show, don't just tell. Providing 1-2 examples of desired input/output pairs vastly improves structure.\n\n"
     "Bad: 'Write a product review.'\n"
     "Good: 'Here are 2 reviews I wrote... now write one in that same voice.'", 
     GREEN),
    ("3. Ask for Explicit Structure", 
     "Request specific output formats like bulleted lists, standard markdown tables, or plain JSON structures.\n\n"
     "Bad: 'Compare these two.'\n"
     "Good: 'Compare A and B in a markdown table with columns for Cost, Speed, and Ease of Use.'", 
     AMBER)
]

for i, (title, desc, color) in enumerate(principles):
    x = M + i * Inches(3.9)
    card = make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.1), Inches(3.5), Inches(4.5), WHITE)
    card.line.color.rgb = color; card.line.width = Pt(2)
    add_textbox(s, x + Inches(0.2), Inches(1.3), Inches(3.1), Inches(0.4), title, 14, color, True)
    add_textbox(s, x + Inches(0.2), Inches(1.8), Inches(3.1), Inches(3.6), desc, 11, BLACK)

# callout
make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.8), Inches(11.3), Inches(0.8), RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(s, Inches(1.3), Inches(5.9), Inches(10.7), Inches(0.6),
    "Pro Tip: If the model gives a poor answer, don't restart. Reply with: 'That's too technical, explain it simpler' or 'Format it as a table.'", 11, GREEN, True)

slide_number(s, 7)
notes(s, "Prompt engineering is 'garbage in, garbage out'. Being specific, providing examples, and asking for structure will improve your outputs by 10x immediately.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — TEMPERATURE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Temperature: The Creativity Knob")

add_textbox(s, M, Inches(1.0), Inches(11.3), Inches(1.0),
    "Temperature is a setting from 0.0 to 1.0+ that controls how 'creative' or 'variable' the model's outputs are.\n"
    "Low temperature forces the model to pick the most mathematically probable words. High temperature introduces randomness.",
    13, BLACK)

# Zones
zones = [
    ("Low Temperature (0.0)", "Highly predictable & consistent.", "Deterministic",
     "• Writing software code\n• Math calculations\n• Fact extraction\n• Structured JSON output", RGBColor(0x1A, 0x2D, 0x4D)),
    
    ("Balanced Temperature (0.5)", "Blend of focus and fluidity.", "Optimal for most tasks",
     "• Summarizing articles\n• Answering general Q&A\n• Drafting emails\n• Document editing", ACCENT),
     
    ("High Temperature (1.0)", "Highly random and diverse.", "Creative & brainstorm",
     "• Generating fiction stories\n• Naming products\n• Brainstorming ideas\n• Creative marketing slogans", GREEN)
]

for i, (title, subtitle, badge, bullets, color) in enumerate(zones):
    x = M + i * Inches(3.9); y = Inches(2.2)
    card = make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(3.3), WHITE)
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    
    add_textbox(s, x + Inches(0.2), y + Inches(0.15), Inches(3.1), Inches(0.3), title, 13, color, True)
    add_textbox(s, x + Inches(0.2), y + Inches(0.48), Inches(3.1), Inches(0.25), subtitle, 10, MDGRAY)
    
    # badge shape
    make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.8), Inches(2.0), Inches(0.25), color, badge, 8, WHITE)
    
    add_textbox(s, x + Inches(0.2), y + Inches(1.2), Inches(3.1), Inches(1.9), bullets, 11, BLACK)

slide_number(s, 8)
notes(s, "Explain temperature simply. It is the probability scale. Low temp = math and code. High temp = naming products and writing poetry. Most chat platforms manage this automatically, but understanding it helps when building systems.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — GETTING STARTED ROADMAP
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, "Getting Started: Your First Month with AI")

# Timeline
roadmap = [
    ("30 Minutes", "Sign up at Claude.ai or ChatGPT.com\n\nAsk the model to explain a hard concept in 3 different ways. Paste a long email and ask for a 2-sentence summary.", ACCENT),
    ("1 Day", "Use AI for one task on your active agenda\n\nDraft a polite response to a difficult email, format a messy text list into a table, or brainstorm 10 ideas.", GREEN),
    ("1 Week", "Introduce structure & prompt helper templates\n\nSet up a 'custom system prompt' to define your preferred output voice, format, and common parameters.", AMBER),
    ("1 Month", "Standardize workflow automation\n\nAutomate a recurring report summary, deploy a custom prompt pipeline, or share template checklists with team members.", RGBColor(0x1A, 0x2D, 0x4D))
]

# Horizontal bar connecting timeline
make_shape(s, MSO_SHAPE.RECTANGLE, M + Inches(1.0), Inches(3.0), Inches(8.5), Pt(3), LTGRAY)

for i, (time_title, details, color) in enumerate(roadmap):
    x = M + i * Inches(2.8)
    
    # Timeline node node
    make_shape(s, MSO_SHAPE.OVAL, x + Inches(1.0), Inches(2.85), Inches(0.35), Inches(0.35), color)
    
    # Text content
    add_textbox(s, x, Inches(1.3), Inches(2.4), Inches(0.4), time_title, 14, color, True, PP_ALIGN.CENTER)
    
    card = make_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.5), Inches(2.45), Inches(2.8), WHITE)
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    add_textbox(s, x + Inches(0.1), Inches(3.6), Inches(2.25), Inches(2.6), details, 10, BLACK)

slide_number(s, 9)
notes(s, "Encourage the user. AI is a habit, not just a tool. Suggest small, daily steps. By starting with 30-minute experiments and graduating to weekly habits, you earn familiarity and speed.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════
s = new_slide(WHITE)
add_textbox(s, M, Inches(1.8), Inches(11.3), Inches(1.0), "Questions & Discussion", 52, ACCENT, True, PP_ALIGN.CENTER)
add_textbox(s, M, Inches(3.0), Inches(11.3), Inches(0.8), "Ask me anything about generative AI, prompting, or getting started", 18, BLACK, False, PP_ALIGN.CENTER)
add_textbox(s, M, Inches(4.5), Inches(11.3), Inches(0.6), "Visit the Playbook: ai-playbook.pages.dev  \u2022  Claude  \u2022  ChatGPT", 14, MDGRAY, False, PP_ALIGN.CENTER)
slide_number(s, 10)
notes(s, "Open floor for Q&A. Focus on practical starting points and address any initial doubts.")

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ai-playbook-beginner.pptx")
prs.save(output_path)
print(f"\u2713 Saved {output_path}")
